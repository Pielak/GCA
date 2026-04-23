#!/usr/bin/env python3
"""Gera apresentação comercial PPTX do GCA — estilo profissional.

Layout 16:9 (widescreen 13.33" x 7.5"). Cada slide tem áreas claramente
separadas: header (logo + título), content (sem sobreposição) e footer
(rodapé com contador). Cores institucionais consistentes com o doc .docx.

Saída: /home/luiz/GCA/docs/GCA_Apresentacao_Comercial.pptx
"""
from __future__ import annotations

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ─── Configuração ──────────────────────────────────────────────────────

LOGO_PATH = Path("/home/luiz/GCA/logogca.png")
OUT_PATH = Path("/home/luiz/GCA/docs/GCA_Apresentacao_Comercial.pptx")

# Paleta institucional (espelha o documento .docx)
NAVY = RGBColor(0x1E, 0x2D, 0x52)         # logo navy escuro
NAVY_DEEP = RGBColor(0x10, 0x1A, 0x33)
VIOLET = RGBColor(0x6D, 0x28, 0xD9)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
SLATE_LIGHT = RGBColor(0x64, 0x74, 0x8B)
SLATE_BG = RGBColor(0xF1, 0xF5, 0xF9)
EMERALD = RGBColor(0x05, 0x96, 0x69)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_700 = RGBColor(0x37, 0x41, 0x51)

# Layout (Inches) — 16:9 = 13.333 x 7.5
SLIDE_W = 13.333
SLIDE_H = 7.5

# Faixas verticais reservadas
HEADER_H = 1.0          # 0.0-1.0
TITLE_TOP = 1.0
TITLE_H = 0.7           # 1.0-1.7
CONTENT_TOP = 1.85      # 1.85-6.95
CONTENT_H = 5.10
FOOTER_TOP = 7.05       # 7.05-7.4

MARGIN_X = 0.7

# Logo no header (canto superior direito)
LOGO_H = 0.55
LOGO_W = 0.82  # proporção 1168/784 ≈ 1.49 → 0.55 × 1.49 ≈ 0.82
LOGO_RIGHT_PAD = 0.4

# ─── Helpers ───────────────────────────────────────────────────────────

def set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left: float, top: float, width: float, height: float,
             *, fill: RGBColor | None = None, line: RGBColor | None = None,
             line_width: float | None = None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    shape.shadow.inherit = False
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        if line_width is not None:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    # Limpar texto default
    shape.text_frame.text = ""
    return shape


def add_text(slide, left: float, top: float, width: float, height: float,
             text: str, *, size: int = 14, bold: bool = False,
             color: RGBColor = SLATE, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, italic: bool = False,
             font_name: str = "Calibri"):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_logo(slide) -> None:
    """Logo no canto superior direito de cada slide (exceto capa)."""
    left = SLIDE_W - LOGO_W - LOGO_RIGHT_PAD
    top = (HEADER_H - LOGO_H) / 2
    slide.shapes.add_picture(
        str(LOGO_PATH), Inches(left), Inches(top),
        height=Inches(LOGO_H),
    )


def add_header_band(slide, kicker: str, title: str) -> None:
    """Faixa superior consistente: kicker (categoria) + título (slide)."""
    # Faixa fina superior (acento)
    add_rect(slide, 0, 0, SLIDE_W, 0.12, fill=NAVY)
    # Kicker — pequeno, uppercase, navy claro
    add_text(slide, MARGIN_X, 0.32, 6.0, 0.30, kicker.upper(),
             size=10, bold=True, color=VIOLET, anchor=MSO_ANCHOR.MIDDLE)
    # Título — grande
    add_text(slide, MARGIN_X, TITLE_TOP, 11.0, TITLE_H, title,
             size=28, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_logo(slide)


def add_footer(slide, page_num: int, total: int) -> None:
    add_rect(slide, 0, FOOTER_TOP + 0.15, SLIDE_W, 0.02, fill=GRAY_400)
    add_text(slide, MARGIN_X, FOOTER_TOP + 0.20, 8.0, 0.20,
             "GCA — Gestão de Codificação Assistida · Apresentação Comercial",
             size=9, color=SLATE_LIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, SLIDE_W - 1.5, FOOTER_TOP + 0.20, 1.0, 0.20,
             f"{page_num} / {total}",
             size=9, color=SLATE_LIGHT, align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)


# ─── Construtores de slides ────────────────────────────────────────────

def slide_capa(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    # Banda lateral esquerda navy
    add_rect(s, 0, 0, 0.4, SLIDE_H, fill=NAVY)
    # Logo grande centralizado horizontalmente, 1/3 do topo
    logo_h = 2.4
    logo_w = logo_h * 1.49
    s.shapes.add_picture(
        str(LOGO_PATH),
        Inches((SLIDE_W - logo_w) / 2), Inches(1.2),
        height=Inches(logo_h),
    )
    # Título principal
    add_text(s, 1.0, 4.10, SLIDE_W - 2.0, 0.7,
             "Gestão de Codificação Assistida",
             size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # Subtítulo
    add_text(s, 1.0, 4.85, SLIDE_W - 2.0, 0.45,
             "Governança de projetos de software assistida por IA",
             size=18, color=SLATE_LIGHT, align=PP_ALIGN.CENTER, italic=True)
    # Tag colorida embaixo
    add_rect(s, (SLIDE_W - 5.0) / 2, 5.55, 5.0, 0.50,
             fill=VIOLET)
    add_text(s, (SLIDE_W - 5.0) / 2, 5.55, 5.0, 0.50,
             "Apresentação Comercial · Versão 1.0",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Rodapé com data/autor
    add_text(s, 1.0, 6.55, SLIDE_W - 2.0, 0.3,
             f"Autor: Luiz Carlos Pielak  ·  {datetime.now().strftime('%B %Y').capitalize()}",
             size=11, color=SLATE_LIGHT, align=PP_ALIGN.CENTER)


def slide_agenda(prs, page, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_header_band(s, "Sumário", "Como esta apresentação está organizada")

    items = [
        ("01", "O que é o GCA",                     "Definição e propósito"),
        ("02", "Dores que o GCA resolve",           "Seis problemas estruturais"),
        ("03", "Diferenciais de mercado",           "Por que somos diferentes"),
        ("04", "Funcionalidades",                   "Áreas Administrativa e de Projeto"),
        ("05", "O que esperar do produto",          "Visão honesta de maturidade"),
        ("06", "Próximas entregas",                 "Roadmap definido"),
    ]
    # 2 colunas x 3 linhas
    col_w = 5.8
    col_gap = 0.4
    row_h = 1.40
    row_gap = 0.20
    start_x_left = MARGIN_X
    start_x_right = MARGIN_X + col_w + col_gap
    start_y = CONTENT_TOP

    for i, (num, title, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = start_x_left if col == 0 else start_x_right
        y = start_y + row * (row_h + row_gap)
        add_rect(s, x, y, col_w, row_h, fill=SLATE_BG, line=GRAY_400, line_width=0.5)
        # Número grande à esquerda
        add_text(s, x + 0.15, y + 0.10, 0.9, row_h - 0.20,
                 num, size=36, bold=True, color=VIOLET,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        # Título + descrição
        add_text(s, x + 1.10, y + 0.20, col_w - 1.20, 0.50,
                 title, size=15, bold=True, color=NAVY)
        add_text(s, x + 1.10, y + 0.70, col_w - 1.20, row_h - 0.80,
                 desc, size=11, color=SLATE_LIGHT)

    add_footer(s, page, total)


def slide_o_que_e(prs, page, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_header_band(s, "01 · Definição", "O que é o GCA")

    # Frase de definição em destaque
    add_rect(s, MARGIN_X, CONTENT_TOP, SLIDE_W - 2 * MARGIN_X, 1.4,
             fill=NAVY)
    add_text(s, MARGIN_X + 0.3, CONTENT_TOP + 0.10,
             SLIDE_W - 2 * MARGIN_X - 0.6, 1.20,
             "Plataforma instalável por cliente para governança de projetos "
             "de software assistida por Inteligência Artificial. Cobre o "
             "ciclo completo: da solicitação à entrega, sob compartimentalização "
             "dura por projeto e auditoria contínua.",
             size=14, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, italic=True)

    # 4 pilares — cards horizontais
    pillars = [
        ("Instalável", "Uma instância por cliente. Dados nunca saem do ambiente.", VIOLET),
        ("Governada", "OCG como fonte única de verdade. Tudo versionado e auditável.", EMERALD),
        ("Compartimentalizada", "Isolamento dura por projeto em DB, IA, storage e auditoria.", AMBER),
        ("Aberta", "IA configurável, papéis canônicos, integração via padrões abertos.", NAVY),
    ]
    card_y = CONTENT_TOP + 1.85
    card_h = 2.40
    card_w = (SLIDE_W - 2 * MARGIN_X - 0.6) / 4
    for i, (title, desc, color) in enumerate(pillars):
        x = MARGIN_X + i * (card_w + 0.2)
        # Card
        add_rect(s, x, card_y, card_w, card_h, fill=WHITE,
                 line=GRAY_400, line_width=0.75)
        # Faixa superior colorida
        add_rect(s, x, card_y, card_w, 0.30, fill=color)
        # Título
        add_text(s, x + 0.15, card_y + 0.45, card_w - 0.30, 0.55,
                 title, size=16, bold=True, color=color)
        # Descrição
        add_text(s, x + 0.15, card_y + 1.05, card_w - 0.30, card_h - 1.15,
                 desc, size=11, color=GRAY_700)

    add_footer(s, page, total)


def slide_dores(prs, page, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_header_band(s, "02 · Diagnóstico", "Seis dores estruturais que o GCA resolve")

    pains = [
        ("Requisitos voláteis",
         "Mudanças sem rastreio quebram contratos a cada sprint."),
        ("Decisões arquiteturais perdidas",
         "Decisões críticas viram conhecimento tácito de quem ficou."),
        ("Documentação obsoleta",
         "Wikis e PDFs envelhecem; ninguém atualiza após o release."),
        ("Compliance reativo",
         "LGPD, SOX, ISO entram só na auditoria — tarde demais."),
        ("Onboarding lento",
         "Devs novos levam semanas para entender o porquê das escolhas."),
        ("IA sem governança",
         "Provedor escolhido por preferência, sem auditoria nem custo claro."),
    ]
    card_w = (SLIDE_W - 2 * MARGIN_X - 0.40) / 3
    card_h = (CONTENT_H - 0.40) / 2
    for i, (title, desc) in enumerate(pains):
        col = i % 3
        row = i // 3
        x = MARGIN_X + col * (card_w + 0.20)
        y = CONTENT_TOP + row * (card_h + 0.20)
        add_rect(s, x, y, card_w, card_h, fill=SLATE_BG,
                 line=GRAY_400, line_width=0.5)
        # Marca vermelha à esquerda
        add_rect(s, x, y, 0.10, card_h, fill=RED)
        add_text(s, x + 0.25, y + 0.20, card_w - 0.40, 0.50,
                 title, size=15, bold=True, color=NAVY)
        add_text(s, x + 0.25, y + 0.85, card_w - 0.40, card_h - 1.00,
                 desc, size=11, color=GRAY_700)

    add_footer(s, page, total)


def slide_diferenciais(prs, page, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_header_band(s, "03 · Diferenciais", "Por que o GCA é diferente do mercado")

    diffs = [
        ("Soberania", "Uma instância por cliente. Sem SaaS multi-tenant, sem dados em nuvem compartilhada."),
        ("Compartimentalização dura", "project_id obrigatório em todo predicado. Sem vazamento entre projetos."),
        ("OCG como fonte de verdade", "Objeto Canônico versionado. Toda decisão, contrato e regra rastreável."),
        ("IA por criticidade", "Roteamento híbrido §6: Ollama local para baixa, Premium para alta. Custo auditável."),
        ("Provenance em todo artefato", "TestSpec, LiveDoc e Module Details registram OCG, LLM, ingestões e prompt hash."),
        ("Governança auditável", "Hash chain em audit_log_global. Integridade verificável offline."),
    ]
    card_w = (SLIDE_W - 2 * MARGIN_X - 0.30) / 2
    card_h = (CONTENT_H - 0.40) / 3
    for i, (title, desc) in enumerate(diffs):
        col = i % 2
        row = i // 2
        x = MARGIN_X + col * (card_w + 0.30)
        y = CONTENT_TOP + row * (card_h + 0.20)
        add_rect(s, x, y, card_w, card_h, fill=WHITE,
                 line=NAVY, line_width=1.0)
        # Bullet circular
        bullet = s.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + (card_h - 0.45) / 2),
            Inches(0.45), Inches(0.45),
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = VIOLET
        bullet.line.fill.background()
        bullet.text_frame.text = ""
        # Título e descrição
        add_text(s, x + 0.85, y + 0.18, card_w - 1.0, 0.45,
                 title, size=14, bold=True, color=NAVY)
        add_text(s, x + 0.85, y + 0.62, card_w - 1.0, card_h - 0.78,
                 desc, size=11, color=GRAY_700)

    add_footer(s, page, total)


def slide_funcionalidades_admin(prs, page, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_header_band(s, "04 · Funcionalidades", "Área Administrativa")

    intro = ("Acessada por Admin (ou Sustentação para subset). Sidebar com nove "
             "entradas administra a instância sem necessidade de operar projetos.")
    add_text(s, MARGIN_X, CONTENT_TOP, SLIDE_W - 2 * MARGIN_X, 0.55,
             intro, size=12, color=SLATE_LIGHT, italic=True)

    items = [
        "Dashboard Global — KPIs cross-projeto",
        "Gestão de Projetos — lifecycle (ativo/pausado/inativo/órfão)",
        "Gestão de Usuários — promover/rebaixar/excluir + convite",
        "Auditoria Global — hash chain integridade verificável",
        "Métricas — uso de IA por projeto, custo, tokens",
        "Backups — agregado cross-projeto + quick action",
        "Incidentes — tickets escalados a Admin/Sustentação",
        "Equipe Sustentação — flag is_support independente",
        "Releases — versionamento, snapshot pré-destrutiva",
    ]
    # 3 colunas x 3 linhas
    grid_top = CONTENT_TOP + 0.75
    card_w = (SLIDE_W - 2 * MARGIN_X - 0.40) / 3
    card_h = (CONTENT_H - 0.95) / 3
    for i, text in enumerate(items):
        col = i % 3
        row = i // 3
        x = MARGIN_X + col * (card_w + 0.20)
        y = grid_top + row * (card_h + 0.15)
        add_rect(s, x, y, card_w, card_h, fill=SLATE_BG,
                 line=GRAY_400, line_width=0.5)
        # Número
        add_text(s, x + 0.15, y + 0.10, 0.55, card_h - 0.20,
                 f"{i+1:02d}", size=22, bold=True, color=VIOLET,
                 anchor=MSO_ANCHOR.MIDDLE)
        # Texto
        if " — " in text:
            title, desc = text.split(" — ", 1)
            add_text(s, x + 0.80, y + 0.12, card_w - 0.95, 0.45,
                     title, size=12, bold=True, color=NAVY)
            add_text(s, x + 0.80, y + 0.55, card_w - 0.95, card_h - 0.60,
                     desc, size=10, color=GRAY_700)
        else:
            add_text(s, x + 0.80, y + 0.20, card_w - 0.95, card_h - 0.30,
                     text, size=12, color=NAVY,
                     anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


def slide_funcionalidades_projeto(prs, page, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_header_band(s, "04 · Funcionalidades", "Área de Projeto (visão do GP)")

    intro = ("Acessada via /p/{slug} por Dev, Tester, QA e GP. Sidebar do "
             "projeto cobre todo o ciclo, agrupada em quatro frentes.")
    add_text(s, MARGIN_X, CONTENT_TOP, SLIDE_W - 2 * MARGIN_X, 0.55,
             intro, size=12, color=SLATE_LIGHT, italic=True)

    # Agrupado em 4 categorias
    groups = [
        ("Descoberta", VIOLET, [
            "Equipe", "OCG", "Repositórios", "Ingestão", "Gatekeeper", "Arguidor",
        ]),
        ("Execução", EMERALD, [
            "CodeGen", "Backlog", "Roadmap", "Deploy Plan",
        ]),
        ("Qualidade", AMBER, [
            "Plano de Testes", "Tester Review", "Doc Viva", "Definition of Done",
        ]),
        ("Operação", NAVY, [
            "Configurações", "Audit", "Backups", "Incidentes", "Métricas",
        ]),
    ]
    grid_top = CONTENT_TOP + 0.75
    col_w = (SLIDE_W - 2 * MARGIN_X - 0.6) / 4
    col_h = CONTENT_H - 0.95
    for i, (gname, color, items) in enumerate(groups):
        x = MARGIN_X + i * (col_w + 0.20)
        # Cabeçalho do grupo
        add_rect(s, x, grid_top, col_w, 0.45, fill=color)
        add_text(s, x, grid_top, col_w, 0.45, gname, size=13, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Lista
        list_y = grid_top + 0.55
        for j, item in enumerate(items):
            iy = list_y + j * 0.42
            add_rect(s, x, iy, col_w, 0.36, fill=SLATE_BG,
                     line=GRAY_400, line_width=0.4)
            # Bullet quadrado pequeno
            add_rect(s, x + 0.12, iy + 0.13, 0.10, 0.10, fill=color)
            add_text(s, x + 0.30, iy, col_w - 0.40, 0.36,
                     item, size=11, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


def slide_funcionalidades_ia(prs, page, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_header_band(s, "04 · Funcionalidades", "Inteligência Artificial e Automação")

    # 4 destaques
    feats = [
        ("OCG — Objeto Canônico de Governança",
         "Pipeline de 8 agentes (analyzer + 7 pilares paralelos + consolidator). "
         "Versionado em ocg_delta_log. Toda decisão arquitetural passa por aqui."),
        ("Roteamento híbrido por criticidade",
         "Baixa criticidade → Ollama local (custo zero). Alta criticidade → "
         "Premium (Anthropic/OpenAI). Auditável por chamada em ai_usage_log."),
        ("Plano de Testes gerado por LLM",
         "Cinco tipos: unit/integration/e2e (Ollama, por módulo) + "
         "security/compliance (Premium, globais). Provenance completa por spec."),
        ("Documentação Viva real",
         "module_doc por módulo (Ollama) + index e architecture consolidados "
         "(Premium). Stale detection automática quando OCG evolui."),
    ]
    card_w = (SLIDE_W - 2 * MARGIN_X - 0.30) / 2
    card_h = (CONTENT_H - 0.30) / 2
    for i, (title, desc) in enumerate(feats):
        col = i % 2
        row = i // 2
        x = MARGIN_X + col * (card_w + 0.30)
        y = CONTENT_TOP + row * (card_h + 0.20)
        add_rect(s, x, y, card_w, card_h, fill=WHITE,
                 line=VIOLET, line_width=1.0)
        # Faixa lateral colorida
        add_rect(s, x, y, 0.18, card_h, fill=VIOLET)
        add_text(s, x + 0.35, y + 0.18, card_w - 0.50, 0.50,
                 title, size=14, bold=True, color=NAVY)
        add_text(s, x + 0.35, y + 0.75, card_w - 0.50, card_h - 0.90,
                 desc, size=11, color=GRAY_700)

    add_footer(s, page, total)


def slide_o_que_esperar(prs, page, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_header_band(s, "05 · Maturidade", "O que esperar do produto — visão honesta")

    cols = [
        ("V1 — Entregue agora", EMERALD, [
            "10 MVPs fechados",
            "1.162 testes de regressão",
            "RBAC canônico de 5 papéis",
            "OCG + 7 pilares Gatekeeper",
            "CodeGen em 7 linguagens",
            "TestSpecs e LiveDocs reativos",
            "Backups por projeto + restore",
            "Releases versionadas com snapshot",
            "Tickets de incidente roteados",
        ]),
        ("V1+ — Próximos releases", AMBER, [
            "SSO via OIDC (Azure AD, Okta, Google)",
            "SAML 2.0 corporativo",
            "Wizard de instalação assistida",
            "Auto-upgrade preservando dados",
            "Multi-instância federada (read-only)",
            "Backup off-site assinado",
        ]),
        ("V2 — Visão de produto", VIOLET, [
            "Federação cross-instância (benchmarks)",
            "Marketplace de prompts curados",
            "Hardening operacional avançado",
            "Versionamento de releases declarativas",
            "Trocas autorizadas entre instâncias",
            "Dashboard executivo cross-instância",
        ]),
    ]
    col_w = (SLIDE_W - 2 * MARGIN_X - 0.40) / 3
    col_h = CONTENT_H - 0.10
    for i, (cname, color, items) in enumerate(cols):
        x = MARGIN_X + i * (col_w + 0.20)
        y = CONTENT_TOP
        # Cabeçalho
        add_rect(s, x, y, col_w, 0.55, fill=color)
        add_text(s, x, y, col_w, 0.55, cname, size=14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Lista
        list_y = y + 0.65
        item_h = (col_h - 0.65 - (len(items) - 1) * 0.08) / len(items)
        for j, item in enumerate(items):
            iy = list_y + j * (item_h + 0.08)
            add_rect(s, x, iy, col_w, item_h, fill=SLATE_BG,
                     line=GRAY_400, line_width=0.4)
            # Marca lateral colorida
            add_rect(s, x, iy, 0.10, item_h, fill=color)
            add_text(s, x + 0.22, iy, col_w - 0.32, item_h,
                     item, size=10, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


def slide_proximas_entregas(prs, page, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_header_band(s, "06 · Roadmap", "Próximas entregas definidas")

    items = [
        ("MVP 11", "Identity Federation — SSO via OIDC",
         "Azure AD, Okta, Google Workspace, Keycloak. JIT provisioning. "
         "Fallback bcrypt sempre preservado.",
         "3-4 dias"),
        ("MVP 12", "SAML 2.0 corporativo",
         "ACS endpoint, metadata exchange, claim mapping. "
         "Para clientes que ainda padronizam SAML.",
         "5-7 dias"),
        ("MVP 13", "Data Federation — métricas cross-instância",
         "Endpoints /federation/* com mTLS, anonimização, allow-list "
         "explícito. Exige emenda formal ao contrato §3.",
         "8-12 dias"),
        ("Hardening", "Wizard de instalação + auto-upgrade",
         "Setup interativo Ubuntu/Windows, smoke tests, releases "
         "declarativas com snapshot pré-destrutiva.",
         "Sub-projeto"),
    ]
    item_h = (CONTENT_H - 0.30) / len(items)
    for i, (tag, title, desc, effort) in enumerate(items):
        y = CONTENT_TOP + i * (item_h + 0.10)
        add_rect(s, MARGIN_X, y, SLIDE_W - 2 * MARGIN_X, item_h - 0.10,
                 fill=WHITE, line=NAVY, line_width=1.0)
        # Tag à esquerda
        add_rect(s, MARGIN_X, y, 1.4, item_h - 0.10, fill=NAVY)
        add_text(s, MARGIN_X, y, 1.4, item_h - 0.10, tag,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Título + desc
        add_text(s, MARGIN_X + 1.6, y + 0.10, SLIDE_W - 2 * MARGIN_X - 3.4, 0.40,
                 title, size=14, bold=True, color=NAVY)
        add_text(s, MARGIN_X + 1.6, y + 0.55, SLIDE_W - 2 * MARGIN_X - 3.4,
                 item_h - 0.70, desc, size=11, color=GRAY_700)
        # Esforço à direita
        eff_x = SLIDE_W - MARGIN_X - 1.6
        add_rect(s, eff_x, y + 0.20, 1.5, item_h - 0.50,
                 fill=SLATE_BG, line=GRAY_400, line_width=0.4)
        add_text(s, eff_x, y + 0.20, 1.5, item_h - 0.50,
                 effort, size=12, bold=True, color=VIOLET,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


def slide_metricas(prs, page, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_header_band(s, "Estado do Produto", "Métricas atuais — abril de 2026")

    metrics = [
        ("10",   "MVPs fechados",         VIOLET),
        ("1.162", "Testes de regressão",   EMERALD),
        ("7",    "Linguagens em CodeGen", AMBER),
        ("5",    "Papéis canônicos",      NAVY),
        ("0",    "Dívidas bloqueantes",   EMERALD),
        ("100%", "Compartimentalização",  VIOLET),
    ]
    card_w = (SLIDE_W - 2 * MARGIN_X - 0.50) / 3
    card_h = (CONTENT_H - 0.30) / 2
    for i, (value, label, color) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = MARGIN_X + col * (card_w + 0.25)
        y = CONTENT_TOP + row * (card_h + 0.20)
        add_rect(s, x, y, card_w, card_h, fill=WHITE,
                 line=color, line_width=1.5)
        # Número grande
        add_text(s, x, y + 0.20, card_w, card_h * 0.55,
                 value, size=54, bold=True, color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Label embaixo
        add_text(s, x, y + card_h * 0.65, card_w, card_h * 0.30,
                 label, size=13, color=SLATE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page, total)


def slide_proximos_passos(prs, page, total):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_header_band(s, "Próximos Passos", "Como avaliar o GCA")

    steps = [
        ("01", "Avaliação técnica",
         "Demo guiada de 60 minutos com pipeline completo (questionário, OCG, "
         "Roadmap, CodeGen, Plano de Testes, Doc Viva)."),
        ("02", "Piloto em projeto real",
         "Instância dedicada com seu projeto atual. Acompanhamento de 2 semanas "
         "para validar adequação, custos de IA e workflow do time."),
        ("03", "Implantação",
         "Instalação assistida em Ubuntu ou Windows, transferência de "
         "conhecimento, documentação operacional, contrato de Sustentação."),
    ]
    step_h = (CONTENT_H - 0.40) / len(steps)
    for i, (num, title, desc) in enumerate(steps):
        y = CONTENT_TOP + i * (step_h + 0.20)
        add_rect(s, MARGIN_X, y, SLIDE_W - 2 * MARGIN_X, step_h,
                 fill=SLATE_BG, line=NAVY, line_width=1.0)
        # Número grande à esquerda
        add_rect(s, MARGIN_X, y, 1.6, step_h, fill=NAVY)
        add_text(s, MARGIN_X, y, 1.6, step_h,
                 num, size=46, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Título + desc
        add_text(s, MARGIN_X + 1.8, y + 0.20, SLIDE_W - 2 * MARGIN_X - 2.0, 0.50,
                 title, size=18, bold=True, color=NAVY)
        add_text(s, MARGIN_X + 1.8, y + 0.80, SLIDE_W - 2 * MARGIN_X - 2.0,
                 step_h - 1.00, desc, size=12, color=GRAY_700)

    add_footer(s, page, total)


def slide_obrigado(prs):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    # Banda lateral esquerda navy (espelha a capa, mantém consistência)
    add_rect(s, 0, 0, 0.4, SLIDE_H, fill=NAVY)
    # Logo grande centralizado (sobre fundo branco — visibilidade ok)
    logo_h = 1.8
    logo_w = logo_h * 1.49
    s.shapes.add_picture(
        str(LOGO_PATH),
        Inches((SLIDE_W - logo_w) / 2), Inches(1.4),
        height=Inches(logo_h),
    )
    # Mensagem
    add_text(s, 1.0, 3.80, SLIDE_W - 2.0, 0.85,
             "Obrigado.",
             size=54, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, 1.0, 4.75, SLIDE_W - 2.0, 0.50,
             "Vamos governar seu próximo projeto juntos.",
             size=18, color=SLATE_LIGHT, align=PP_ALIGN.CENTER, italic=True)
    # Contato
    add_rect(s, (SLIDE_W - 6.0) / 2, 5.65, 6.0, 0.60, fill=VIOLET)
    add_text(s, (SLIDE_W - 6.0) / 2, 5.65, 6.0, 0.60,
             "Luiz Carlos Pielak  ·  pielak.ctba@gmail.com",
             size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Linha sutil de assinatura
    add_text(s, 1.0, 6.55, SLIDE_W - 2.0, 0.30,
             "GCA — Gestão de Codificação Assistida",
             size=10, color=SLATE_LIGHT, align=PP_ALIGN.CENTER)


# ─── Build ─────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # Slides na ordem (capa + 9 conteúdo + obrigado)
    builders = [
        slide_capa,
        slide_agenda,
        slide_o_que_e,
        slide_dores,
        slide_diferenciais,
        slide_funcionalidades_admin,
        slide_funcionalidades_projeto,
        slide_funcionalidades_ia,
        slide_o_que_esperar,
        slide_proximas_entregas,
        slide_metricas,
        slide_proximos_passos,
        slide_obrigado,
    ]
    total = len(builders)
    # Capa = 1 (sem footer); demais com footer indicando page/total
    for i, fn in enumerate(builders):
        if fn in (slide_capa, slide_obrigado):
            fn(prs)
        else:
            fn(prs, i + 1, total)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"✓ Apresentação gerada: {OUT_PATH}")
    print(f"  Slides: {total}  Tamanho: {OUT_PATH.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
