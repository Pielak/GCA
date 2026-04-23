#!/usr/bin/env python3
"""Gera apresentação institucional do GCA (Gestão de Codificação Assistida).

Destino: sócio. Tom: direto, sucinto, demonstração de capacidade.
Inclui screenshots reais da área administrativa e de gestão de projeto
+ quadro comparativo GCA × concorrentes.

Saída:
  - docs/GCA_Apresentacao_Socio.pptx (intermediário, pra edição futura)
  - docs/GCA_Apresentacao_Socio.pdf   (entregável final)

Conversão PPTX→PDF via LibreOffice headless.
"""
from __future__ import annotations

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# ─── Configuração ──────────────────────────────────────────────────────

LOGO_PATH = Path("/home/luiz/GCA/logogca.png")
SHOTS = Path("/home/luiz/GCA/docs/screenshots")
OUT_PPTX = Path("/home/luiz/GCA/docs/GCA_Apresentacao_Socio.pptx")
OUT_PDF = Path("/home/luiz/GCA/docs/GCA_Apresentacao_Socio.pdf")

# Landscape 16:9 — padrão presentation.
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.5
CONTENT_W = SLIDE_W - 2 * MARGIN

# Paleta Observatory (dark theme).
NAVY_DEEP = RGBColor(0x0A, 0x0F, 0x1F)
SURFACE_1 = RGBColor(0x10, 0x17, 0x2A)
SURFACE_2 = RGBColor(0x17, 0x20, 0x39)
SURFACE_3 = RGBColor(0x1E, 0x29, 0x3B)
EDGE = RGBColor(0x2B, 0x36, 0x4D)
EDGE_STRONG = RGBColor(0x3D, 0x4A, 0x66)

VIOLET = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_SOFT = RGBColor(0xA7, 0x8B, 0xFA)
VIOLET_GLOW = RGBColor(0x4C, 0x1D, 0x95)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
SKY = RGBColor(0x38, 0xBD, 0xF8)

INK_HIGH = RGBColor(0xF1, 0xF5, 0xF9)
INK_MED = RGBColor(0xCB, 0xD5, 0xE1)
INK_LOW = RGBColor(0x94, 0xA3, 0xB8)
INK_DIM = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xF8, 0xFA, 0xFC)


# ─── Helpers ───────────────────────────────────────────────────────────


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _set_line(shape, color: RGBColor | None, width: float = 0.75) -> None:
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_rect(slide, x, y, w, h, fill=SURFACE_1, line=None, corner=0.15):
    assert x >= 0 and y >= 0 and x + w <= SLIDE_W + 0.01 and y + h <= SLIDE_H + 0.01
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    _set_fill(shape, fill)
    _set_line(shape, line, 0.75)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *,
             size=13, color=INK_HIGH, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Inter", italic=False, line_spacing=1.2):
    assert x >= 0 and y >= 0 and x + w <= SLIDE_W + 0.01 and y + h <= SLIDE_H + 0.01
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_pill(slide, x, y, text, *, fill=SURFACE_3, text_color=INK_MED,
             size=9, border=None, char_w=0.075, padding_x=0.2):
    approx_w = max(0.6, char_w * len(text) + padding_x * 2)
    h = 0.32
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
        Inches(approx_w), Inches(h),
    )
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    _set_fill(shape, fill)
    _set_line(shape, border, 0.5)
    shape.shadow.inherit = False
    tb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(approx_w), Inches(h)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Inter"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    return approx_w


def paint_background(slide, color=NAVY_DEEP) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_W), Inches(SLIDE_H),
    )
    _set_fill(bg, color)
    _set_line(bg, None)
    bg.shadow.inherit = False


def add_header_accent(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_W), Inches(0.09),
    )
    _set_fill(shape, VIOLET)
    _set_line(shape, None)
    shape.shadow.inherit = False


def add_footer(slide, idx, total) -> None:
    y = SLIDE_H - 0.42
    add_text(
        slide, MARGIN, y, CONTENT_W * 0.6, 0.3,
        "GCA — Gestão de Codificação Assistida  ·  equipe de desenvolvimento",
        size=8, color=INK_DIM, bold=True,
    )
    add_text(
        slide, MARGIN + CONTENT_W * 0.6, y, CONTENT_W * 0.4, 0.3,
        f"{idx:02d} / {total:02d}",
        size=8, color=INK_LOW, bold=True, align=PP_ALIGN.RIGHT,
    )


def add_title(slide, kicker, title) -> None:
    add_text(slide, MARGIN, 0.28, CONTENT_W, 0.4,
             kicker, size=10, color=VIOLET_SOFT, bold=True)
    add_text(slide, MARGIN, 0.65, CONTENT_W, 0.95,
             title, size=26, color=INK_HIGH, bold=True, line_spacing=1.05)


def add_logo_mark(slide, x, y, w=1.8, h=0.55) -> None:
    add_rect(slide, x, y, w, h, fill=PAPER, line=EDGE_STRONG, corner=0.3)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Inches(x + 0.12), Inches(y + 0.08),
            height=Inches(h - 0.16),
        )


def add_screenshot_frame(slide, x, y, max_w, max_h, path, caption=None) -> None:
    add_rect(slide, x, y, max_w, max_h, fill=SURFACE_2, line=EDGE_STRONG)
    add_rect(slide, x, y, max_w, 0.32, fill=SURFACE_3, line=None, corner=0.12)
    for i, color in enumerate([RED, AMBER, EMERALD]):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + 0.14 + i * 0.2), Inches(y + 0.1),
            Inches(0.12), Inches(0.12),
        )
        _set_fill(dot, color)
        _set_line(dot, None)
        dot.shadow.inherit = False
    add_text(slide, x + 0.9, y + 0.03, max_w - 1.2, 0.28,
             caption or "gca.local", size=8, color=INK_LOW,
             font="JetBrains Mono")
    inner_x = x + 0.08
    inner_y = y + 0.4
    inner_w = max_w - 0.16
    inner_h = max_h - 0.48
    if path.exists():
        pic = slide.shapes.add_picture(
            str(path), Inches(inner_x), Inches(inner_y),
            width=Inches(inner_w),
        )
        actual_h_in = pic.height / 914400
        if actual_h_in > inner_h:
            pic.height = Inches(inner_h)
            pic.width = Inches(inner_h * (pic.width / pic.height))


# ─── Slides ────────────────────────────────────────────────────────────


def slide_01_capa(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)

    # Glow decorativo
    glow = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5),
        Inches(5.5), Inches(5.5),
    )
    _set_fill(glow, VIOLET_GLOW)
    _set_line(glow, None)
    glow.shadow.inherit = False

    add_logo_mark(slide, MARGIN, 0.5, w=2.4, h=0.8)

    add_pill(slide, MARGIN, 1.7,
             "APRESENTAÇÃO INTERNA · CAPACIDADES E FUNCIONALIDADES",
             fill=VIOLET_GLOW, text_color=VIOLET_SOFT,
             size=10, border=VIOLET, char_w=0.088)

    add_text(slide, MARGIN, 2.4, CONTENT_W * 0.85, 2.8,
             "GCA\nGestão de Codificação\nAssistida.",
             size=52, color=INK_HIGH, bold=True, line_spacing=1.02)

    add_text(slide, MARGIN, 5.3, CONTENT_W * 0.85, 1.3,
             "Plataforma instalável on-premises que governa o ciclo completo de "
             "desenvolvimento assistido por IA — do requisito ao código, "
             "com auditoria, rastreabilidade e integração ao ecossistema "
             "corporativo do cliente.",
             size=14, color=INK_MED, line_spacing=1.4)

    add_text(slide, MARGIN, SLIDE_H - 0.95, CONTENT_W, 0.35,
             "Preparado por equipe de desenvolvimento GCA",
             size=11, color=INK_LOW, bold=True)


def slide_02_sintese(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "SÍNTESE",
              "O que o GCA é, o que resolve, como entrega.")

    cols = 3
    gap = 0.3
    cw = (CONTENT_W - gap * (cols - 1)) / cols
    ch = 5.2
    cy = 1.75

    cards = [
        ("O que é",
         "Produto instalável por cliente — uma instância por empresa, "
         "sem SaaS compartilhado.",
         "Pipeline canônico de 10 etapas do questionário até o release, "
         "cada etapa lendo e atualizando o OCG (fonte única de verdade "
         "do projeto).",
         "IA é configurável por projeto (Anthropic, OpenAI, Gemini, Ollama "
         "local) — cliente escolhe provedor e fallback.",
         VIOLET),
        ("O que resolve",
         "Requisitos dispersos em PDF/e-mail/WhatsApp — vão pro OCG vivo.",
         "IA que gera código sem contrato e vira dívida técnica — GCA "
         "gera dentro de contrato + rastreabilidade.",
         "Documentação morta no dia seguinte — LiveDocs + ERS IEEE 830 "
         "regenerados automaticamente.",
         "Auditoria reativa — SHA-256 hash chain em toda ação.",
         EMERALD),
        ("Como entrega",
         "8 linguagens de codegen (Python, Node×2, Java×2, Kotlin, C#, "
         "Go, PHP, C++).",
         "ERS vivo no Git do projeto (docs/ERS.md) com matriz de "
         "rastreabilidade § 4 IEEE 830.",
         "Integração com Jira, Trello, Sonar, Snyk, gitleaks, Slack, "
         "Microsoft Teams — via adapter pattern canônico.",
         "Governança: 5 papéis RBAC, LGPD/GDPR, quarentena de PII, "
         "billing por IA, backups automáticos.",
         SKY),
    ]

    cx = MARGIN
    for title, *bullets, accent in cards:
        add_rect(slide, cx, cy, cw, ch, fill=SURFACE_1, line=EDGE)
        add_rect(slide, cx, cy, cw, 0.1, fill=accent, corner=0)
        add_text(slide, cx + 0.3, cy + 0.3, cw - 0.6, 0.5,
                 title, size=18, color=INK_HIGH, bold=True)
        by = cy + 0.95
        for b in bullets:
            add_text(slide, cx + 0.3, by, 0.3, 0.5,
                     "•", size=14, color=accent, bold=True)
            add_text(slide, cx + 0.55, by, cw - 0.85, 0.95,
                     b, size=11, color=INK_MED, line_spacing=1.35)
            by += 1.0
        cx += cw + gap

    add_footer(slide, idx, total)


def slide_03_pipeline(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "PIPELINE CANÔNICO",
              "Do questionário à documentação viva, em 10 etapas governadas.")

    add_text(slide, MARGIN, 1.65, CONTENT_W, 0.4,
             "Cada etapa lê o OCG antes, atualiza depois. Auditoria SHA-256 "
             "em toda transição. Humano decide, IA executa.",
             size=12, color=INK_LOW, italic=True, line_spacing=1.3)

    stages = [
        ("1", "Questionário", "49 perguntas externas", VIOLET),
        ("2", "Repos + Ingestão", "Código legado + docs", VIOLET),
        ("3", "Gatekeeper", "7 pilares + score", VIOLET),
        ("4", "Arguidor", "Perguntas dirigidas", VIOLET),
        ("5", "OCG", "Consolidação", EMERALD),
        ("6", "CodeGen", "Scaffold + módulos", EMERALD),
        ("7", "Testes", "Unit + Int + E2E", EMERALD),
        ("8", "QA / Revisão", "Aprovação humana", EMERALD),
        ("9", "Roadmap", "Fases priorizadas", SKY),
        ("10", "Docs + ERS", "IEEE 830 atualizado", SKY),
    ]

    cols = 5
    rows = 2
    gap_x, gap_y = 0.2, 0.25
    cw = (CONTENT_W - gap_x * (cols - 1)) / cols
    ch = 1.85
    start_y = 2.3

    for i, (num, title, sub, color) in enumerate(stages):
        r = i // cols
        c = i % cols
        cx = MARGIN + c * (cw + gap_x)
        cy = start_y + r * (ch + gap_y)
        add_rect(slide, cx, cy, cw, ch, fill=SURFACE_1, line=EDGE)
        add_rect(slide, cx + 0.15, cy + 0.15, 0.5, 0.45,
                 fill=color, corner=0.35)
        add_text(slide, cx + 0.15, cy + 0.15, 0.5, 0.45,
                 num, size=12, color=INK_HIGH, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx + 0.15, cy + 0.7, cw - 0.3, 0.45,
                 title, size=12, color=INK_HIGH, bold=True)
        add_text(slide, cx + 0.15, cy + 1.15, cw - 0.3, 0.6,
                 sub, size=9, color=INK_LOW, line_spacing=1.25)

    cy = start_y + 2 * ch + gap_y + 0.1
    add_text(slide, MARGIN, cy, CONTENT_W, 0.5,
             "Sem OCG válido, o pipeline não avança. Sem aprovação do GP, "
             "não há merge. Governança é embutida, não opcional.",
             size=14, color=VIOLET_SOFT, bold=True, italic=True,
             align=PP_ALIGN.CENTER)

    add_footer(slide, idx, total)


def slide_04_gatekeeper(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "GATEKEEPER · AVALIAÇÃO FORMAL",
              "7 pilares com score + bloqueio automático.")

    add_screenshot_frame(slide, MARGIN, 1.8, 8.5, 4.9,
                          SHOTS / "gatekeeper.png",
                          caption="/projects/:id/gatekeeper")

    rx = MARGIN + 8.5 + 0.25
    rw = CONTENT_W - 8.5 - 0.25

    add_text(slide, rx, 1.9, rw, 0.45,
             "O semáforo antes do código.",
             size=14, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.4, rw, 1.7,
             "O Gatekeeper calcula score global a partir dos 7 pilares do "
             "OCG (Negócio, Compliance, Escopo, Performance, Arquitetura, "
             "Dados, Segurança). Cada pilar com peso canônico e blockers "
             "explícitos.",
             size=10, color=INK_MED, line_spacing=1.4)

    bx = rx
    by = 4.25
    add_rect(slide, bx, by, rw, 2.35, fill=SURFACE_1, line=EDGE)
    add_text(slide, bx + 0.2, by + 0.15, rw - 0.4, 0.3,
             "REGRAS DE APROVAÇÃO", size=8, color=VIOLET_SOFT, bold=True)
    rules = [
        ("P7 < 70", "BLOCKED (Segurança)", RED),
        ("P2 < 70", "BLOCKED (Compliance)", RED),
        ("≥ 90", "READY", EMERALD),
        ("≥ 75", "NEEDS_REVIEW", AMBER),
        ("< 75", "AT_RISK", AMBER),
    ]
    ry = by + 0.5
    for cond, status, color in rules:
        add_text(slide, bx + 0.25, ry, rw * 0.35, 0.3,
                 cond, size=10, color=INK_MED, font="JetBrains Mono", bold=True)
        add_text(slide, bx + 0.25 + rw * 0.35, ry, rw * 0.6, 0.3,
                 f"→ {status}", size=10, color=color, bold=True)
        ry += 0.32

    add_footer(slide, idx, total)


def slide_05_codegen(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "GERAÇÃO DE CÓDIGO",
              "Scaffold determinístico + IA dentro do contrato do OCG.")

    add_screenshot_frame(slide, MARGIN, 1.75, 8.3, 3.8,
                          SHOTS / "codegen.png",
                          caption="/projects/:id/codegen")

    # 8 linguagens em linha horizontal
    ly = 5.65
    add_text(slide, MARGIN, ly, CONTENT_W, 0.32,
             "8 LINGUAGENS CANÔNICAS · 9 SCAFFOLDERS DETERMINÍSTICOS",
             size=10, color=VIOLET_SOFT, bold=True)
    langs = [
        "Python", "Node.js / Express", "Node.js / NestJS",
        "Java / Spring", "Java / Quarkus", "Kotlin / Spring",
        "C# / ASP.NET", "Go", "PHP / Laravel", "C++ / CMake",
    ]
    lx = MARGIN
    lyy = ly + 0.4
    for name in langs:
        w = add_pill(slide, lx, lyy, name, fill=SURFACE_2,
                     text_color=INK_HIGH, border=EDGE_STRONG, size=9)
        lx += w + 0.1
        if lx > MARGIN + CONTENT_W - 1:
            lx = MARGIN
            lyy += 0.42

    rx = MARGIN + 8.3 + 0.25
    rw = CONTENT_W - 8.3 - 0.25
    add_text(slide, rx, 1.8, rw, 0.4,
             "Código, testes e docstrings juntos.",
             size=13, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.3, rw, 1.2,
             "Cada arquivo gerado é rastreável ao requisito RF/RNF/BR via "
             "matriz §4 do ERS + trilha de auditoria SHA-256.",
             size=10, color=INK_MED, line_spacing=1.35)

    ey = 3.6
    add_rect(slide, rx, ey, rw, 2.0, fill=SURFACE_3, line=VIOLET)
    add_text(slide, rx + 0.2, ey + 0.12, rw - 0.4, 0.35,
             "EXPANSÍVEL", size=9, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx + 0.2, ey + 0.45, rw - 0.4, 0.45,
             "Nova linguagem em ~2-3 dias.",
             size=13, color=INK_HIGH, bold=True)
    add_text(slide, rx + 0.2, ey + 0.95, rw - 0.4, 1.0,
             "Rust, Swift, Scala, Cobol, saídas no-code (Bubble, Retool, "
             "n8n) entram via novo scaffolder — mesmo padrão canônico.",
             size=10, color=INK_MED, line_spacing=1.35)

    add_footer(slide, idx, total)


def slide_06_roadmap(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "ROADMAP DO PROJETO",
              "Backlog organizado em fases executáveis.")

    add_screenshot_frame(slide, MARGIN, 1.75, 8.5, 4.9,
                          SHOTS / "roadmap.png",
                          caption="/projects/:id/roadmap")

    rx = MARGIN + 8.5 + 0.25
    rw = CONTENT_W - 8.5 - 0.25
    add_text(slide, rx, 1.9, rw, 0.45,
             "Do questionário ao deploy.",
             size=14, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.4, rw, 1.2,
             "O Roadmap consolida módulos em fases canônicas — Fundação, "
             "Funcionalidades, Complementos — com filtros por categoria e "
             "status.",
             size=10, color=INK_MED, line_spacing=1.35)

    items = [
        ("Fases", "ordem de deploy sugerida"),
        ("Filtros", "por tipo e prioridade"),
        ("Próxima ação", "recomendada pelo sistema"),
        ("Compartimentalizado", "por projeto — sem vazamento"),
    ]
    iy = 3.75
    for t, sub in items:
        add_rect(slide, rx, iy, rw, 0.65, fill=SURFACE_1, line=EDGE)
        add_text(slide, rx + 0.2, iy + 0.08, rw - 0.4, 0.3,
                 t, size=11, color=INK_HIGH, bold=True)
        add_text(slide, rx + 0.2, iy + 0.35, rw - 0.4, 0.3,
                 sub, size=9, color=INK_LOW)
        iy += 0.7

    add_footer(slide, idx, total)


def slide_07_docs(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "DOCUMENTAÇÃO VIVA + ERS IEEE 830",
              "Docs que se atualizam a cada evento do pipeline.")

    add_screenshot_frame(slide, MARGIN, 1.75, 8.3, 4.9,
                          SHOTS / "docs_vivas.png",
                          caption="/projects/:id/docs")

    rx = MARGIN + 8.3 + 0.25
    rw = CONTENT_W - 8.3 - 0.25
    add_text(slide, rx, 1.9, rw, 0.45,
             "Arquitetura · Módulos · ERS.",
             size=13, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.4, rw, 1.3,
             "Cada doc registra procedência: versão do OCG, ingestões, "
             "provedor e modelo de IA usado.",
             size=10, color=INK_MED, line_spacing=1.35)

    hy = 3.85
    add_rect(slide, rx, hy, rw, 2.7, fill=SURFACE_3, line=VIOLET)
    add_text(slide, rx + 0.2, hy + 0.15, rw - 0.4, 0.3,
             "ERS IEEE 830 VIVO", size=8, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx + 0.2, hy + 0.5, rw - 0.4, 0.5,
             "docs/ERS.md no Git do projeto.",
             size=12, color=INK_HIGH, bold=True)
    add_text(slide, rx + 0.2, hy + 1.05, rw - 0.4, 1.5,
             "Regenerado automaticamente em 4 seções canônicas:\n"
             "  • Introdução + glossário vivo\n"
             "  • Descrição geral\n"
             "  • Requisitos (RF/RNF/BR categorizados)\n"
             "  • Matriz de rastreabilidade",
             size=10, color=INK_MED, line_spacing=1.4)

    add_footer(slide, idx, total)


def slide_08_provedores(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "PROVEDOR DE IA À ESCOLHA",
              "Cliente escolhe. Fallback pra local. Custos visíveis.")

    add_screenshot_frame(slide, MARGIN, 1.75, 8.3, 4.2,
                          SHOTS / "provedores.png",
                          caption="/projects/:id/settings · Provedor de IA")

    add_text(slide, MARGIN, 6.05, 8.3, 0.5,
             "Multi-provedor por projeto. Fallback hierárquico. Zero vendor "
             "lock-in.",
             size=11, color=INK_LOW, italic=True, line_spacing=1.3)

    rx = MARGIN + 8.3 + 0.25
    rw = CONTENT_W - 8.3 - 0.25

    add_text(slide, rx, 1.85, rw, 0.45,
             "Nem refém, nem cego.",
             size=14, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.35, rw, 1.2,
             "Use o modelo que melhor atende seu objetivo — qualidade, "
             "custo, latência ou privacidade.",
             size=10, color=INK_MED, line_spacing=1.35)

    benefits = [
        ("Premium", "decisão crítica (OCG, arquitetura)", VIOLET),
        ("Local / Ollama", "lote repetitivo (docs de módulo)", EMERALD),
        ("Fallback", "se provedor cair, outro assume", SKY),
        ("Vault", "chaves nunca em plaintext", AMBER),
    ]
    by = 3.75
    for t, sub, color in benefits:
        add_rect(slide, rx, by, rw, 0.75, fill=SURFACE_1, line=EDGE)
        add_rect(slide, rx, by, 0.1, 0.75, fill=color, corner=0)
        add_text(slide, rx + 0.25, by + 0.1, rw - 0.4, 0.3,
                 t, size=11, color=INK_HIGH, bold=True)
        add_text(slide, rx + 0.25, by + 0.4, rw - 0.4, 0.3,
                 sub, size=9, color=INK_LOW)
        by += 0.8

    add_footer(slide, idx, total)


def slide_09_metricas(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "MÉTRICAS · CUSTO DE IA VISÍVEL",
              "Controle total do gasto por projeto e operação.")

    add_screenshot_frame(slide, MARGIN, 1.75, 8.5, 4.9,
                          SHOTS / "metricas.png",
                          caption="/projects/:id/metrics")

    rx = MARGIN + 8.5 + 0.25
    rw = CONTENT_W - 8.5 - 0.25
    add_text(slide, rx, 1.9, rw, 0.45,
             "Sem surpresa na fatura.",
             size=14, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.4, rw, 1.15,
             "Cada chamada de LLM registra provedor, modelo, operação, "
             "tokens in/out e custo em USD.",
             size=10, color=INK_MED, line_spacing=1.35)

    hy = 3.7
    add_rect(slide, rx, hy, rw, 2.8, fill=SURFACE_1, line=EDGE)
    add_text(slide, rx + 0.2, hy + 0.12, rw - 0.4, 0.3,
             "VISÍVEL EM TEMPO REAL", size=8, color=VIOLET_SOFT, bold=True)
    items = [
        ("Chamadas LLM totais", "contador acumulado"),
        ("Tokens in / out", "separados"),
        ("Custo estimado (USD)", "por combinação"),
        ("Provider × operação", "breakdown detalhado"),
        ("Janela ajustável", "24h · 7d · 30d"),
    ]
    iy = hy + 0.45
    for label, value in items:
        add_text(slide, rx + 0.25, iy, rw * 0.55, 0.3,
                 label, size=10, color=INK_MED)
        add_text(slide, rx + 0.25 + rw * 0.55, iy, rw * 0.4, 0.3,
                 value, size=10, color=INK_HIGH, bold=True,
                 align=PP_ALIGN.RIGHT, font="JetBrains Mono")
        iy += 0.38

    add_footer(slide, idx, total)


def slide_10_backups(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "BACKUPS AUTOMÁTICOS",
              "Cron diário + sob demanda + rollback em um clique.")

    add_screenshot_frame(slide, MARGIN, 1.75, 8.3, 4.4,
                          SHOTS / "backup.png",
                          caption="/projects/:id/backups")

    rx = MARGIN + 8.3 + 0.25
    rw = CONTENT_W - 8.3 - 0.25
    add_text(slide, rx, 1.85, rw, 0.45,
             "Compartimentalizado por projeto.",
             size=13, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.35, rw, 1.3,
             "Mantém os últimos 10 snapshots. Rollback restaura OCG, "
             "módulos, test_specs, LiveDocs, glossário e configs.",
             size=10, color=INK_MED, line_spacing=1.35)

    points = [
        ("Horário", "00 12 * * * (parametrizável)"),
        ("Retenção", "últimos 10 snapshots"),
        ("Integridade", "hash SHA por backup"),
        ("Manual", "GP ou Admin via botão"),
        ("Catch-up", "auto se servidor estava down"),
    ]
    py = 3.8
    for t, sub in points:
        add_rect(slide, rx, py, rw, 0.5, fill=SURFACE_1, line=EDGE)
        add_text(slide, rx + 0.2, py + 0.08, rw * 0.35, 0.3,
                 t, size=10, color=INK_HIGH, bold=True)
        add_text(slide, rx + 0.2 + rw * 0.35, py + 0.08, rw * 0.6, 0.3,
                 sub, size=9, color=INK_LOW)
        py += 0.55

    add_footer(slide, idx, total)


def slide_11_seguranca(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "SEGURANÇA & COMPLIANCE EMBARCADOS",
              "LGPD, GDPR e aderência ISO 27001 — por design, não como add-on.")

    # 2 colunas: Compliance legal · Controles de segurança
    col_w = (CONTENT_W - 0.3) / 2
    col_h = 5.2
    cy = 1.8

    # Coluna esquerda — Compliance legal
    add_rect(slide, MARGIN, cy, col_w, col_h, fill=SURFACE_1, line=EMERALD)
    add_rect(slide, MARGIN, cy, col_w, 0.1, fill=EMERALD, corner=0)
    add_text(slide, MARGIN + 0.3, cy + 0.3, col_w - 0.6, 0.5,
             "Compliance legal",
             size=20, color=INK_HIGH, bold=True)
    add_text(slide, MARGIN + 0.3, cy + 0.85, col_w - 0.6, 0.4,
             "LGPD · GDPR — atendimento real",
             size=11, color=EMERALD, bold=True, italic=True)

    compliance = [
        ("Quarentena automática de PII",
         "CPF, CNPJ, cartão, telefone, email detectados em ingestão e "
         "retidos até decisão do GP."),
        ("Trilha de consentimento",
         "Todo processamento registrado com base legal, timestamp e "
         "actor — exportável para ANPD."),
        ("Direito de esquecimento",
         "Exclusão em cascata preserva a cadeia de auditoria sem "
         "expor PII removida."),
        ("Residência de dados",
         "On-premises por design — dados nunca saem da infraestrutura "
         "do cliente."),
    ]
    iy = cy + 1.4
    for t, sub in compliance:
        add_text(slide, MARGIN + 0.35, iy, 0.3, 0.3,
                 "✓", size=14, color=EMERALD, bold=True)
        add_text(slide, MARGIN + 0.7, iy, col_w - 1.0, 0.32,
                 t, size=11, color=INK_HIGH, bold=True)
        add_text(slide, MARGIN + 0.7, iy + 0.3, col_w - 1.0, 0.55,
                 sub, size=9, color=INK_MED, line_spacing=1.3)
        iy += 0.92

    # Coluna direita — Controles técnicos / ISO 27001 aderente
    rx = MARGIN + col_w + 0.3
    rw = col_w
    add_rect(slide, rx, cy, rw, col_h, fill=SURFACE_1, line=VIOLET)
    add_rect(slide, rx, cy, rw, 0.1, fill=VIOLET, corner=0)
    add_text(slide, rx + 0.3, cy + 0.3, rw - 0.6, 0.5,
             "Controles técnicos",
             size=20, color=INK_HIGH, bold=True)
    add_text(slide, rx + 0.3, cy + 0.85, rw - 0.6, 0.4,
             "Aderente a ISO/IEC 27001 · Anexo A",
             size=11, color=VIOLET_SOFT, bold=True, italic=True)

    controls = [
        ("Hash chain SHA-256",
         "Auditoria encadeada: cada evento carrega hash do anterior; "
         "qualquer alteração quebra a cadeia."),
        ("Vault pgcrypto",
         "Chaves de IA, tokens Git, credenciais de integração sempre "
         "encrypted; nunca em plaintext."),
        ("RBAC canônico de 5 papéis",
         "Admin · GP · Dev · Tester · QA — fronteira de ação clara por "
         "papel, enforced no backend."),
        ("Compartimentalização por projeto",
         "Todo dado filtrado por project_id; zero vazamento cross-tenant "
         "mesmo na instância única."),
    ]
    iy = cy + 1.4
    for t, sub in controls:
        add_text(slide, rx + 0.35, iy, 0.3, 0.3,
                 "✓", size=14, color=VIOLET_SOFT, bold=True)
        add_text(slide, rx + 0.7, iy, rw - 1.0, 0.32,
                 t, size=11, color=INK_HIGH, bold=True)
        add_text(slide, rx + 0.7, iy + 0.3, rw - 1.0, 0.55,
                 sub, size=9, color=INK_MED, line_spacing=1.3)
        iy += 0.92

    # Rodapé honesto sobre certificação
    fy = cy + col_h + 0.15
    add_text(slide, MARGIN, fy, CONTENT_W, 0.3,
             "ISO 27001: GCA oferece controles técnicos compatíveis com o Anexo A (A.5 a A.14). "
             "Certificação formal é da instância operada pelo cliente.",
             size=9, color=INK_DIM, italic=True, align=PP_ALIGN.CENTER)

    add_footer(slide, idx, total)


def slide_12_integracoes(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "INTEGRAÇÕES EXTERNAS",
              "GCA consome ferramentas que o cliente já tem.")

    # Hoje
    lx, ly, lw, lh = MARGIN, 1.85, 5.9, 4.9
    add_rect(slide, lx, ly, lw, lh, fill=SURFACE_1, line=EMERALD)
    add_rect(slide, lx, ly, lw, 0.1, fill=EMERALD, corner=0)
    add_text(slide, lx + 0.3, ly + 0.25, lw - 0.6, 0.4,
             "DISPONÍVEL HOJE", size=10, color=EMERALD, bold=True)
    add_text(slide, lx + 0.3, ly + 0.7, lw - 0.6, 0.6,
             "Ecossistema corporativo", size=18, color=INK_HIGH, bold=True)

    available = [
        ("GitHub · GitLab · Bitbucket", "commits, PRs/MRs, webhooks"),
        ("Jira · Trello", "backlog ↔ issue tracker"),
        ("Sonar · Snyk · gitleaks", "findings alimentam P7 do OCG"),
        ("Slack · Microsoft Teams", "eventos canônicos no canal"),
    ]
    ry = ly + 1.55
    for name, sub in available:
        add_rect(slide, lx + 0.3, ry, lw - 0.6, 0.7, fill=SURFACE_2, line=EDGE)
        ch = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(lx + 0.45), Inches(ry + 0.2),
            Inches(0.3), Inches(0.3),
        )
        _set_fill(ch, EMERALD)
        _set_line(ch, None)
        ch.shadow.inherit = False
        add_text(slide, lx + 0.45, ry + 0.2, 0.3, 0.3,
                 "✓", size=11, color=NAVY_DEEP, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, lx + 0.85, ry + 0.12, lw - 1.1, 0.3,
                 name, size=11, color=INK_HIGH, bold=True)
        add_text(slide, lx + 0.85, ry + 0.4, lw - 1.1, 0.3,
                 sub, size=9, color=INK_LOW)
        ry += 0.8

    # Sob demanda
    rx = MARGIN + lw + 0.25
    rw = CONTENT_W - lw - 0.25
    add_rect(slide, rx, ly, rw, lh, fill=SURFACE_1, line=VIOLET)
    add_rect(slide, rx, ly, rw, 0.1, fill=VIOLET, corner=0)
    add_text(slide, rx + 0.3, ly + 0.25, rw - 0.6, 0.4,
             "SOB DEMANDA · ~2-4 DIAS CADA",
             size=10, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx + 0.3, ly + 0.7, rw - 0.6, 0.6,
             "Padrão adapter canônico",
             size=18, color=INK_HIGH, bold=True)

    future = [
        ("Linear · Asana · Monday", "issue trackers modernos"),
        ("Mattermost · Discord", "chat alternativo"),
        ("ChatOps bi-direcional", "aprovar via botão no canal"),
        ("Figma MCP", "design system enterprise"),
    ]
    ry = ly + 1.55
    for name, sub in future:
        add_rect(slide, rx + 0.3, ry, rw - 0.6, 0.7, fill=SURFACE_2, line=EDGE)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(rx + 0.45), Inches(ry + 0.2),
            Inches(0.3), Inches(0.3),
        )
        _set_fill(dot, VIOLET)
        _set_line(dot, None)
        dot.shadow.inherit = False
        add_text(slide, rx + 0.45, ry + 0.2, 0.3, 0.3,
                 "+", size=14, color=INK_HIGH, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, rx + 0.85, ry + 0.12, rw - 1.1, 0.3,
                 name, size=11, color=INK_HIGH, bold=True)
        add_text(slide, rx + 0.85, ry + 0.4, rw - 1.1, 0.3,
                 sub, size=9, color=INK_LOW)
        ry += 0.8

    add_footer(slide, idx, total)


def slide_13_comparativo(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "GCA × CONCORRENTES DIRETOS",
              "Quadro comparativo — o que só o GCA entrega.")

    tools = ["Copilot", "Cursor", "v0/Lovable", "GCA"]
    rows_data = [
        ("Autocomplete de IA", True, True, True, True),
        ("Geração de módulo completo", False, "parcial", True, True),
        ("Contexto vivo do projeto (OCG)", False, False, False, True),
        ("Classificação IEEE 830 (RF/RNF/BR)", False, False, False, True),
        ("ERS auto-regenerado no Git", False, False, False, True),
        ("Matriz requisito × teste × código", False, False, False, True),
        ("RBAC + auditoria SHA-256", False, False, False, True),
        ("LGPD / GDPR + quarentena PII", False, False, False, True),
        ("Multi-linguagem governada", False, "parcial", False, "8 langs"),
        ("Controle de custo de IA por projeto", False, False, False, True),
        ("Integração Jira/Sonar/Slack nativa", False, False, False, True),
        ("Instalável on-premises", False, False, False, True),
    ]

    tx = MARGIN
    ty = 1.75
    tw = CONTENT_W
    row_h = 0.33
    cell_feat_w = tw * 0.44
    cell_tool_w = (tw - cell_feat_w) / len(tools)

    # Header
    add_rect(slide, tx, ty, tw, row_h, fill=SURFACE_2, line=EDGE)
    add_text(slide, tx + 0.15, ty, cell_feat_w - 0.15, row_h,
             "Recurso", size=10, color=INK_LOW, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    for i, t in enumerate(tools):
        cx = tx + cell_feat_w + i * cell_tool_w
        is_gca = t == "GCA"
        if is_gca:
            add_rect(slide, cx, ty, cell_tool_w, row_h,
                     fill=VIOLET_GLOW, line=VIOLET)
        add_text(slide, cx, ty, cell_tool_w, row_h,
                 t, size=10, color=VIOLET_SOFT if is_gca else INK_HIGH,
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ty += row_h

    for i, row in enumerate(rows_data):
        feat, *vals = row
        bg = SURFACE_1 if i % 2 == 0 else NAVY_DEEP
        add_rect(slide, tx, ty, tw, row_h, fill=bg, line=EDGE)
        add_text(slide, tx + 0.15, ty, cell_feat_w - 0.15, row_h,
                 feat, size=9, color=INK_MED, anchor=MSO_ANCHOR.MIDDLE)
        for j, v in enumerate(vals):
            is_gca = j == len(vals) - 1
            cx = tx + cell_feat_w + j * cell_tool_w
            if is_gca:
                add_rect(slide, cx, ty, cell_tool_w, row_h,
                         fill=VIOLET_GLOW, line=None)
            if v is True:
                text, color, size = "✓", EMERALD, 12
            elif v is False:
                text, color, size = "✗", RED, 12
            else:
                text, color, size = str(v), AMBER, 9
            add_text(slide, cx, ty, cell_tool_w, row_h,
                     text, size=size, color=color, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ty += row_h

    # Callout
    cy = ty + 0.2
    add_text(slide, MARGIN, cy, CONTENT_W, 0.55,
             "Copilot/Cursor aceleram o teclado. GCA governa a entrega.",
             size=14, color=VIOLET_SOFT, bold=True, italic=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, MARGIN, cy + 0.6, CONTENT_W, 0.5,
             "GCA não compete com Sonar, Jira ou Slack — integra com eles sob governança.",
             size=11, color=INK_LOW, italic=True, align=PP_ALIGN.CENTER)

    add_footer(slide, idx, total)


def slide_14_investimento(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "INVESTIMENTO COMPARATIVO",
              "Pricing Enterprise dos concorrentes × GCA (consultivo).")

    # Tabela comparativa de preços
    # Coluna 1: Ferramenta · 2: Plano Enterprise · 3: Preço público
    tx = MARGIN
    ty = 1.75
    tw = CONTENT_W
    row_h = 0.42
    c1 = tw * 0.28  # Ferramenta
    c2 = tw * 0.42  # Plano / escopo
    c3 = tw * 0.30  # Preço

    # Header
    add_rect(slide, tx, ty, tw, row_h, fill=SURFACE_2, line=EDGE)
    add_text(slide, tx + 0.2, ty, c1 - 0.2, row_h,
             "Ferramenta", size=10, color=INK_LOW, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, tx + c1, ty, c2, row_h,
             "Plano Enterprise", size=10, color=INK_LOW, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT)
    add_text(slide, tx + c1 + c2, ty, c3, row_h,
             "Preço público de referência", size=10, color=INK_LOW, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    ty += row_h

    competitors = [
        ("GitHub Copilot Enterprise",
         "Assistente de código + chat no repo",
         "US$ 39 / usuário / mês"),
        ("Cursor Business",
         "IDE com IA + contexto de repo",
         "US$ 40 / usuário / mês"),
        ("GitHub Advanced Security",
         "SAST + secret scanning + supply chain",
         "US$ 49 / usuário / mês"),
        ("SonarQube Enterprise",
         "Self-hosted, análise estática corporativa",
         "US$ 20k+ / ano (até 1M LOC)"),
        ("Snyk Enterprise",
         "SCA + SAST + IaC + container scanning",
         "US$ 98 / dev / mês (custom)"),
        ("Atlassian Jira Enterprise",
         "Gestão de projetos + workflows avançados",
         "US$ 90 / usuário / mês"),
        ("Slack Enterprise Grid",
         "Canal corporativo + compliance + audit logs",
         "US$ 15-30 / usuário / mês"),
    ]

    for i, (tool, plan, price) in enumerate(competitors):
        bg = SURFACE_1 if i % 2 == 0 else NAVY_DEEP
        add_rect(slide, tx, ty, tw, row_h, fill=bg, line=EDGE)
        add_text(slide, tx + 0.2, ty, c1 - 0.2, row_h,
                 tool, size=10, color=INK_HIGH, bold=True,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + c1, ty, c2, row_h,
                 plan, size=9, color=INK_MED,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, tx + c1 + c2, ty, c3, row_h,
                 price, size=10, color=INK_MED, bold=True,
                 font="JetBrains Mono",
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        ty += row_h

    # Linha GCA destacada
    add_rect(slide, tx, ty, tw, row_h * 1.15, fill=VIOLET_GLOW, line=VIOLET)
    add_text(slide, tx + 0.2, ty, c1 - 0.2, row_h * 1.15,
             "GCA", size=13, color=VIOLET_SOFT, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, tx + c1, ty, c2, row_h * 1.15,
             "Plataforma instalável completa",
             size=10, color=INK_HIGH, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, tx + c1 + c2, ty, c3, row_h * 1.15,
             "Negociável conforme demanda",
             size=11, color=VIOLET_SOFT, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    ty += row_h * 1.15 + 0.15

    # Callout / interpretação
    add_text(slide, MARGIN, ty, CONTENT_W, 0.5,
             "Um stack corporativo completo combinando as ferramentas acima passa facilmente de US$ 250 / dev / mês — sem incluir integração entre elas.",
             size=11, color=INK_MED, italic=True, align=PP_ALIGN.CENTER)

    # Disclaimer
    add_text(slide, MARGIN, SLIDE_H - 0.8, CONTENT_W, 0.3,
             "Valores de referência pública (2025). Cotações variam por plano, volume, duração do contrato e negociação direta com fornecedor.",
             size=8, color=INK_DIM, italic=True, align=PP_ALIGN.CENTER)

    add_footer(slide, idx, total)


def slide_15_fechamento(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)

    # Glow decorativo
    glow = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-2), Inches(SLIDE_H - 3),
        Inches(6), Inches(6),
    )
    _set_fill(glow, VIOLET_GLOW)
    _set_line(glow, None)
    glow.shadow.inherit = False

    add_logo_mark(slide, MARGIN, 0.5, w=2.2, h=0.8)

    add_pill(slide, MARGIN, 1.75,
             "EM OPERAÇÃO · INSTALÁVEL · ON-PREMISES",
             fill=VIOLET_GLOW, text_color=VIOLET_SOFT,
             size=11, border=VIOLET, char_w=0.085)

    add_text(slide, MARGIN, 2.35, CONTENT_W, 3.3,
             "Um produto que trata\ngovernança de IA como\ninfraestrutura, "
             "não processo.",
             size=40, color=INK_HIGH, bold=True, line_spacing=1.04)

    # Stats
    y = 6.0
    stats = [
        ("22", "MVPs entregues"),
        ("8", "linguagens de codegen"),
        ("4", "integrações externas"),
        ("LGPD", "+ ISO 27001 aderente"),
    ]
    sw = (CONTENT_W - 0.45) / 4
    sx = MARGIN
    for num, label in stats:
        add_rect(slide, sx, y, sw, 1.1, fill=SURFACE_1, line=EDGE)
        add_text(slide, sx, y + 0.15, sw, 0.55,
                 num, size=28, color=VIOLET_SOFT, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, sx, y + 0.7, sw, 0.35,
                 label, size=10, color=INK_LOW, bold=True,
                 align=PP_ALIGN.CENTER)
        sx += sw + 0.15

    add_text(slide, MARGIN, SLIDE_H - 0.6, CONTENT_W, 0.35,
             "Equipe de desenvolvimento GCA",
             size=10, color=INK_LOW, bold=True, align=PP_ALIGN.CENTER)


# ─── Main ──────────────────────────────────────────────────────────────


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    builders = [
        slide_01_capa,
        slide_02_sintese,
        slide_03_pipeline,
        slide_04_gatekeeper,
        slide_05_codegen,
        slide_06_roadmap,
        slide_07_docs,
        slide_08_provedores,
        slide_09_metricas,
        slide_10_backups,
        slide_11_seguranca,
        slide_12_integracoes,
        slide_13_comparativo,
        slide_14_investimento,
        slide_15_fechamento,
    ]
    total = len(builders)
    for i, build in enumerate(builders, start=1):
        if build.__name__ == "slide_01_capa":
            build(prs, total)
        else:
            build(prs, i, total)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"✓ PPTX: {OUT_PPTX}  ({OUT_PPTX.stat().st_size / 1024:.1f} KiB)")

    # Converte para PDF via LibreOffice headless.
    result = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(OUT_PDF.parent), str(OUT_PPTX)],
        capture_output=True, text=True, timeout=120,
    )
    if result.returncode != 0:
        print(f"! LibreOffice erro: {result.stderr}")
        return
    if OUT_PDF.exists():
        print(f"✓ PDF:  {OUT_PDF}  ({OUT_PDF.stat().st_size / 1024:.1f} KiB)")
    else:
        print(f"! PDF não foi criado em {OUT_PDF}")


if __name__ == "__main__":
    main()
