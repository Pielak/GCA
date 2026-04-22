#!/usr/bin/env python3
"""Gera apresentação PPTX landscape (16:9) do GCA para LinkedIn/publicação.

16 slides, paleta Observatory, logo em painel de contraste, screenshots
reais da UI (Gatekeeper, Roadmap, CodeGen, Docs Vivas, Backup, Métricas,
Provedores), glossário de acrônimos e mensagem explícita de que o GCA
aumenta performance — não substitui humano.

Saída: /home/luiz/GCA/docs/GCA_LinkedIn_Landscape.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ─── Configuração ──────────────────────────────────────────────────────

LOGO_PATH = Path("/home/luiz/GCA/logogca.png")
SHOTS = Path("/home/luiz/GCA/docs/screenshots")
OUT_PATH = Path("/home/luiz/GCA/docs/GCA_LinkedIn_Landscape.pptx")

# Landscape 16:9 — padrão LinkedIn para documentos PDF e carrossel em imagem.
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.5
CONTENT_W = SLIDE_W - 2 * MARGIN  # 12.333
SAFE_BOTTOM = 7.05

# Paleta Observatory (dark theme).
NAVY_DEEP = RGBColor(0x0A, 0x0F, 0x1F)
SURFACE_1 = RGBColor(0x10, 0x17, 0x2A)
SURFACE_2 = RGBColor(0x17, 0x20, 0x39)
SURFACE_3 = RGBColor(0x1E, 0x29, 0x3B)
EDGE = RGBColor(0x2B, 0x36, 0x4D)
EDGE_STRONG = RGBColor(0x3D, 0x4A, 0x66)

VIOLET = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_SOFT = RGBColor(0xA7, 0x8B, 0xFA)
VIOLET_GLOW = RGBColor(0x4C, 0x1D, 0x95)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
SKY = RGBColor(0x38, 0xBD, 0xF8)
ROSE = RGBColor(0xF4, 0x72, 0xB6)

INK_HIGH = RGBColor(0xF1, 0xF5, 0xF9)
INK_MED = RGBColor(0xCB, 0xD5, 0xE1)
INK_LOW = RGBColor(0x94, 0xA3, 0xB8)
INK_DIM = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAPER = RGBColor(0xF8, 0xFA, 0xFC)  # painel claro pro logo

# ─── Helpers ───────────────────────────────────────────────────────────


def _set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _set_line(shape, color: RGBColor | None, width: float = 0.75) -> None:
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_rect(
    slide, x: float, y: float, w: float, h: float,
    fill: RGBColor = SURFACE_1, line: RGBColor | None = None,
    corner: float = 0.15,
):
    assert x >= 0 and y >= 0 and x + w <= SLIDE_W + 0.01 and y + h <= SLIDE_H + 0.01, (
        f"shape fora do slide: ({x:.2f},{y:.2f} {w:.2f}x{h:.2f})"
    )
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    try:
        shape.adjustments[0] = corner
    except Exception:
        pass
    _set_fill(shape, fill)
    _set_line(shape, line, 0.75)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide, x: float, y: float, w: float, h: float,
    text: str, *,
    size: float = 13, color: RGBColor = INK_HIGH, bold: bool = False,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font: str = "Inter",
    italic: bool = False, line_spacing: float = 1.2,
):
    assert x >= 0 and y >= 0 and x + w <= SLIDE_W + 0.01 and y + h <= SLIDE_H + 0.01, (
        f"textbox fora do slide: ({x:.2f},{y:.2f} {w:.2f}x{h:.2f})"
    )
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_pill(
    slide, x: float, y: float, text: str, *,
    fill: RGBColor = SURFACE_3, text_color: RGBColor = INK_MED,
    size: float = 9, border: RGBColor | None = None,
    char_w: float = 0.072, padding_x: float = 0.2,
):
    approx_w = max(0.6, char_w * len(text) + padding_x * 2)
    h = 0.32
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
        Inches(approx_w), Inches(h),
    )
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    _set_fill(shape, fill)
    _set_line(shape, border, 0.5)
    shape.shadow.inherit = False
    tb = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(approx_w), Inches(h)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Inter"
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    return approx_w


def paint_background(slide, color: RGBColor = NAVY_DEEP) -> None:
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_W), Inches(SLIDE_H),
    )
    _set_fill(bg, color)
    _set_line(bg, None)
    bg.shadow.inherit = False


def add_header_accent(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(SLIDE_W), Inches(0.09),
    )
    _set_fill(shape, VIOLET)
    _set_line(shape, None)
    shape.shadow.inherit = False


def add_footer(slide, idx: int, total: int) -> None:
    y = SLIDE_H - 0.42
    add_text(
        slide, MARGIN, y, CONTENT_W * 0.6, 0.3,
        "GCA — Gestão de Codificação Assistida",
        size=8, color=INK_DIM, bold=True,
    )
    add_text(
        slide, MARGIN + CONTENT_W * 0.6, y, CONTENT_W * 0.4, 0.3,
        f"{idx:02d} / {total:02d}  ·  linkedin.com/in/luizpielak",
        size=8, color=INK_LOW, bold=True, align=PP_ALIGN.RIGHT,
    )


def add_title(slide, kicker: str, title: str) -> None:
    add_text(
        slide, MARGIN, 0.28, CONTENT_W, 0.4,
        kicker,
        size=10, color=VIOLET_SOFT, bold=True,
    )
    add_text(
        slide, MARGIN, 0.65, CONTENT_W, 0.95,
        title,
        size=26, color=INK_HIGH, bold=True, line_spacing=1.05,
    )


def add_logo_mark(slide, x: float, y: float, w: float = 1.8, h: float = 0.55) -> None:
    """Logo em pílula com fundo claro (contraste com tema escuro)."""
    add_rect(slide, x, y, w, h, fill=PAPER, line=EDGE_STRONG, corner=0.3)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Inches(x + 0.12), Inches(y + 0.08),
            height=Inches(h - 0.16),
        )


def add_screenshot_frame(
    slide, x: float, y: float, max_w: float, max_h: float, path: Path,
    caption: str | None = None,
) -> None:
    """Moldura para screenshot — barra de janela simulada + imagem escalada."""
    frame_border = add_rect(slide, x, y, max_w, max_h, fill=SURFACE_2, line=EDGE_STRONG)
    # Topbar (0.35 de altura)
    add_rect(slide, x, y, max_w, 0.32, fill=SURFACE_3, line=None, corner=0.12)
    # Três "dots" da janela
    for i, color in enumerate([RED, AMBER, EMERALD]):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + 0.14 + i * 0.2), Inches(y + 0.1),
            Inches(0.12), Inches(0.12),
        )
        _set_fill(dot, color)
        _set_line(dot, None)
        dot.shadow.inherit = False
    # Path simulado
    add_text(
        slide, x + 0.9, y + 0.03, max_w - 1.2, 0.28,
        caption or "gca.local",
        size=8, color=INK_LOW, font="JetBrains Mono",
    )
    # Imagem ocupa espaço interno abaixo da topbar
    inner_x = x + 0.08
    inner_y = y + 0.4
    inner_w = max_w - 0.16
    inner_h = max_h - 0.48
    # add_picture preserva aspect; ancoramos no canto e deixamos width calcular.
    if path.exists():
        pic = slide.shapes.add_picture(
            str(path), Inches(inner_x), Inches(inner_y),
            width=Inches(inner_w),
        )
        # Se altura real ultrapassar o espaço, reajusta por altura
        actual_h_in = pic.height / 914400
        if actual_h_in > inner_h:
            pic.height = Inches(inner_h)
            pic.width = Inches(inner_h * (pic.width / pic.height))


# ─── Slides ────────────────────────────────────────────────────────────


def slide_01_capa(prs, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)

    # Glow decorativo
    glow = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5),
        Inches(5.5), Inches(5.5),
    )
    _set_fill(glow, VIOLET_GLOW)
    _set_line(glow, None)
    glow.shadow.inherit = False

    # Logo em painel claro (contraste)
    add_logo_mark(slide, MARGIN, 0.45, w=2.3, h=0.75)

    # Pill contextual
    add_pill(
        slide, MARGIN, 1.55,
        "PRODUTO INSTALÁVEL · ON-PREMISES · PT-BR",
        fill=VIOLET_GLOW, text_color=VIOLET_SOFT,
        size=10, border=VIOLET, char_w=0.085,
    )

    # Título principal
    add_text(
        slide, MARGIN, 2.15, CONTENT_W * 0.85, 2.8,
        "Entregar software\ncom IA virou roleta.",
        size=52, color=INK_HIGH, bold=True, line_spacing=1.02,
    )

    # Subtítulo
    add_text(
        slide, MARGIN, 4.6, CONTENT_W * 0.85, 0.9,
        "O GCA é o fim da roleta.",
        size=28, color=VIOLET_SOFT, bold=True,
    )

    # Pitch
    add_text(
        slide, MARGIN, 5.55, CONTENT_W * 0.7, 1.2,
        "Governança de codificação assistida por IA — um produto instalável "
        "que transforma requisitos, arquitetura, geração de código, testes e "
        "documentação em um pipeline auditável ponta a ponta.",
        size=13, color=INK_MED, line_spacing=1.4,
    )

    # Swipe hint
    add_text(
        slide, MARGIN, SLIDE_H - 0.95, CONTENT_W, 0.35,
        "Arraste para ver como  →",
        size=11, color=INK_LOW, bold=True,
    )

    add_footer(slide, 1, total)


def slide_02_glossario(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "ANTES DE COMEÇAR", "Um glossário rápido.")

    entries = [
        ("GCA", "Gestão de Codificação Assistida — o produto."),
        ("OCG", "Objeto de Contexto Global — fonte única de verdade viva do projeto."),
        ("ERS", "Especificação de Requisitos de Software (IEEE 830-1998)."),
        ("RBAC", "Role-Based Access Control — controle de acesso por papel."),
        ("LGPD / GDPR", "Leis de proteção de dados (Brasil / União Europeia)."),
        ("PII", "Personally Identifiable Information — dados pessoais sensíveis."),
        ("LLM", "Large Language Model — modelo de linguagem de IA (Claude, GPT, Llama…)."),
        ("GP", "Gerente de Projeto — papel que aprova, rejeita e libera regeneração."),
        ("RF · RNF · BR", "Requisitos Funcionais · Não-Funcionais · Regras de Negócio."),
        ("CI/CD", "Continuous Integration / Continuous Delivery."),
        ("PR / MR", "Pull Request / Merge Request — proposta de mudança no Git."),
        ("MCP", "Model Context Protocol — padrão de integração de ferramentas de IA."),
    ]

    # Grid 3 colunas × 4 linhas
    cols, rows = 3, 4
    gap_x, gap_y = 0.25, 0.2
    cw = (CONTENT_W - gap_x * (cols - 1)) / cols
    ch = 1.05
    start_y = 1.75
    for i, (sigla, defn) in enumerate(entries):
        r = i // cols
        c = i % cols
        cx = MARGIN + c * (cw + gap_x)
        cy = start_y + r * (ch + gap_y)
        add_rect(slide, cx, cy, cw, ch, fill=SURFACE_1, line=EDGE)
        add_text(
            slide, cx + 0.2, cy + 0.1, cw - 0.4, 0.4,
            sigla, size=13, color=VIOLET_SOFT, bold=True, font="JetBrains Mono",
        )
        add_text(
            slide, cx + 0.2, cy + 0.5, cw - 0.4, 0.55,
            defn, size=9, color=INK_MED, line_spacing=1.25,
        )

    add_footer(slide, idx, total)


def slide_03_dor(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "A DOR", "Você reconhece.")

    pains = [
        ("Requisitos dispersos",
         "PDF, e-mail, WhatsApp, ata de reunião.\nNinguém sabe qual é a versão canônica."),
        ("IA vira dívida técnica",
         "Copilot e Cursor aceleram o teclado,\nmas o código nasce sem contrato."),
        ("Documentação morre no dia seguinte",
         "ERS, arquitetura, manual — tudo escrito\numa vez, desatualizado na sprint seguinte."),
        ("Auditoria é arqueologia",
         "Quem aprovou? Por quê? Sob qual stack?\nNinguém responde sem vasculhar o Slack."),
    ]

    cols = 2
    gap = 0.3
    cw = (CONTENT_W - gap) / cols
    ch = 2.3
    start_y = 1.85
    for i, (t, body) in enumerate(pains):
        cx = MARGIN + (i % cols) * (cw + gap)
        cy = start_y + (i // cols) * (ch + 0.25)
        add_rect(slide, cx, cy, cw, ch, fill=SURFACE_1, line=EDGE)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx + 0.3), Inches(cy + 0.35),
            Inches(0.35), Inches(0.35),
        )
        _set_fill(dot, RED)
        _set_line(dot, None)
        dot.shadow.inherit = False
        add_text(slide, cx + 0.85, cy + 0.3, cw - 1.0, 0.5,
                 t, size=16, color=INK_HIGH, bold=True)
        add_text(slide, cx + 0.85, cy + 0.85, cw - 1.0, 1.35,
                 body, size=12, color=INK_MED, line_spacing=1.35)

    # Conclusão
    add_text(
        slide, MARGIN, 6.5, CONTENT_W, 0.5,
        "Cada sprint, a mesma roleta.",
        size=15, color=VIOLET_SOFT, bold=True, italic=True,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide, idx, total)


def slide_04_nao_substitui(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(slide, "POSICIONAMENTO", "O GCA não substitui o ser humano.")

    # Frase de destaque
    add_text(
        slide, MARGIN, 1.8, CONTENT_W, 1.3,
        "Ele aumenta performance e qualidade.\nA decisão continua humana.",
        size=24, color=VIOLET_SOFT, bold=True, line_spacing=1.15,
    )

    # Duas colunas: "IA faz" vs "Humano decide"
    col_w = (CONTENT_W - 0.4) / 2
    col_h = 3.2
    cy = 3.4

    # Coluna IA
    add_rect(slide, MARGIN, cy, col_w, col_h, fill=SURFACE_1, line=EDGE)
    add_rect(slide, MARGIN, cy, col_w, 0.1, fill=SKY, corner=0)
    add_text(slide, MARGIN + 0.3, cy + 0.3, col_w - 0.6, 0.5,
             "A IA executa",
             size=18, color=INK_HIGH, bold=True)
    add_text(slide, MARGIN + 0.3, cy + 0.8, col_w - 0.6, 0.4,
             "rascunhos, consolidações, varreduras",
             size=11, color=SKY, italic=True, bold=True)
    ia_items = [
        "Gera OCG a partir do questionário e ingestões",
        "Sugere stack, pilares, achados críticos",
        "Produz scaffold de código, specs de teste, docs",
        "Extrai glossário e regras de negócio",
        "Calcula custos, gera relatórios, rastreia",
    ]
    ly = cy + 1.4
    for item in ia_items:
        add_text(slide, MARGIN + 0.45, ly, 0.25, 0.35,
                 "•", size=13, color=SKY, bold=True)
        add_text(slide, MARGIN + 0.75, ly, col_w - 1.0, 0.35,
                 item, size=11, color=INK_MED)
        ly += 0.33

    # Coluna Humano
    hx = MARGIN + col_w + 0.4
    add_rect(slide, hx, cy, col_w, col_h, fill=SURFACE_1, line=EDGE)
    add_rect(slide, hx, cy, col_w, 0.1, fill=EMERALD, corner=0)
    add_text(slide, hx + 0.3, cy + 0.3, col_w - 0.6, 0.5,
             "O humano decide",
             size=18, color=INK_HIGH, bold=True)
    add_text(slide, hx + 0.3, cy + 0.8, col_w - 0.6, 0.4,
             "aprovação, prioridade, liberação",
             size=11, color=EMERALD, italic=True, bold=True)
    humano_items = [
        "GP aprova OCG, módulos, glossário, regeneração do ERS",
        "Dev revisa código gerado, edita, commita",
        "Tester aprova specs de teste, executa, registra",
        "QA revisa execução, libera ou bloqueia",
        "Admin governa chaves, políticas e papéis",
    ]
    ly = cy + 1.4
    for item in humano_items:
        add_text(slide, hx + 0.45, ly, 0.25, 0.35,
                 "✓", size=13, color=EMERALD, bold=True)
        add_text(slide, hx + 0.75, ly, col_w - 1.0, 0.35,
                 item, size=11, color=INK_MED)
        ly += 0.33

    add_footer(slide, idx, total)


def slide_05_ocg(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "COMO FUNCIONA · 01",
        "OCG — a inteligência viva do projeto.",
    )

    # Layout split: esquerda mockup, direita explicação
    left_x, left_y, left_w, left_h = MARGIN, 1.8, 7.3, 5.0
    right_x = left_x + left_w + 0.3
    right_w = CONTENT_W - left_w - 0.3

    # Mockup
    add_rect(slide, left_x, left_y, left_w, left_h, fill=SURFACE_1, line=EDGE)
    add_rect(slide, left_x, left_y, left_w, 0.45, fill=SURFACE_3, corner=0.12)
    add_text(slide, left_x + 0.25, left_y + 0.1, left_w - 0.5, 0.3,
             "gca.local  ·  /projects/financehub-pro/ocg",
             size=9, color=INK_LOW, font="JetBrains Mono")

    # Header do OCG
    add_text(slide, left_x + 0.25, left_y + 0.6, left_w - 0.5, 0.45,
             "OCG v7  ·  FinanceHub Pro",
             size=16, color=INK_HIGH, bold=True)
    add_pill(slide, left_x + left_w - 2.3, left_y + 0.67,
             "NEEDS_REVIEW · 82/100",
             fill=AMBER, text_color=NAVY_DEEP, size=9,
             border=None, char_w=0.08)
    add_text(slide, left_x + 0.25, left_y + 1.05, left_w - 0.5, 0.3,
             "atualizado há 12 min  ·  delta: DOCUMENT_INGESTED",
             size=9, color=INK_LOW, italic=True)

    # Pilares
    pillars = [
        ("P1", 88, EMERALD), ("P2", 74, AMBER), ("P3", 91, EMERALD),
        ("P4", 79, AMBER), ("P5", 86, EMERALD), ("P6", 81, EMERALD),
        ("P7", 72, AMBER),
    ]
    py = left_y + 1.5
    pw = (left_w - 0.5) / 7
    px = left_x + 0.25
    for name, score, color in pillars:
        add_rect(slide, px + 0.04, py, pw - 0.08, 0.95, fill=SURFACE_2, line=EDGE)
        add_text(slide, px + 0.04, py + 0.1, pw - 0.08, 0.28,
                 name, size=10, color=INK_LOW, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, px + 0.04, py + 0.36, pw - 0.08, 0.4,
                 str(score), size=17, color=color, bold=True, align=PP_ALIGN.CENTER)
        bw = (pw - 0.3) * (score / 100)
        add_rect(slide, px + 0.15, py + 0.78, pw - 0.3, 0.07,
                 fill=EDGE, line=None, corner=0.5)
        if bw > 0.02:
            add_rect(slide, px + 0.15, py + 0.78, bw, 0.07,
                     fill=color, line=None, corner=0.5)
        px += pw

    # Stack
    sy = left_y + 2.7
    add_text(slide, left_x + 0.25, sy, left_w - 0.5, 0.25,
             "STACK RECOMENDADO",
             size=8, color=INK_DIM, bold=True)
    add_text(slide, left_x + 0.25, sy + 0.28, left_w - 0.5, 0.4,
             "Python · FastAPI · PostgreSQL 16 · Redis · React + Vite",
             size=11, color=INK_HIGH, bold=True, font="JetBrains Mono")

    # Critical finding
    fy = left_y + 3.5
    add_rect(slide, left_x + 0.25, fy, left_w - 0.5, 0.75, fill=SURFACE_3, line=AMBER)
    add_text(slide, left_x + 0.45, fy + 0.1, 0.5, 0.3,
             "⚠", size=16, color=AMBER, bold=True)
    add_text(slide, left_x + 0.9, fy + 0.08, left_w - 1.4, 0.32,
             "CRITICAL_FINDING · P2 LGPD",
             size=10, color=AMBER, bold=True)
    add_text(slide, left_x + 0.9, fy + 0.4, left_w - 1.4, 0.3,
             "Ausência de política de retenção bloqueia CodeGen.",
             size=9, color=INK_MED)

    # Delta log
    dy = left_y + 4.4
    add_text(slide, left_x + 0.25, dy, left_w - 0.5, 0.22,
             "ÚLTIMOS DELTAS",
             size=8, color=INK_DIM, bold=True)
    deltas = [
        ("+", "BUSINESS_RULES: Nota fiscal até 24h", EMERALD),
        ("~", "PILLAR_SCORES.P2: 81 → 74", AMBER),
    ]
    ddy = dy + 0.25
    for sign, text, color in deltas:
        add_text(slide, left_x + 0.25, ddy, 0.25, 0.25,
                 sign, size=11, color=color, bold=True, font="JetBrains Mono")
        add_text(slide, left_x + 0.55, ddy, left_w - 0.8, 0.25,
                 text, size=9, color=INK_MED, font="JetBrains Mono")
        ddy += 0.25

    # Lado direito — explicação
    ry = 1.9
    bullets = [
        ("12 seções canônicas", "perfil, pilares, stack, compliance, riscos…"),
        ("Versionado a cada evento", "QUESTIONNAIRE_APPROVED · DOCUMENT_INGESTED…"),
        ("Expande com bom dado", "mais contexto, maior confiança."),
        ("Contrai com conflito", "reduz scores, marca lacuna, pede resolução."),
        ("Bloqueia pipeline", "P2 ou P7 < 70 → BLOCKED · nada avança."),
    ]
    for t, sub in bullets:
        add_rect(slide, right_x, ry, right_w, 0.85, fill=SURFACE_1, line=EDGE)
        add_text(slide, right_x + 0.2, ry + 0.12, right_w - 0.4, 0.35,
                 t, size=12, color=INK_HIGH, bold=True)
        add_text(slide, right_x + 0.2, ry + 0.45, right_w - 0.4, 0.35,
                 sub, size=10, color=INK_LOW, line_spacing=1.25)
        ry += 0.95

    add_footer(slide, idx, total)


def slide_06_gatekeeper(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "COMO FUNCIONA · 02",
        "Gatekeeper — scoring formal dos 7 pilares.",
    )

    # Screenshot à esquerda (grande)
    add_screenshot_frame(
        slide, MARGIN, 1.8, 8.5, 4.9,
        SHOTS / "gatekeeper.png",
        caption="/projects/:id/gatekeeper",
    )

    # Explicação à direita
    rx = MARGIN + 8.5 + 0.25
    rw = CONTENT_W - 8.5 - 0.25

    add_text(slide, rx, 1.9, rw, 0.45,
             "O semáforo antes do código.",
             size=14, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.4, rw, 1.5,
             "O Gatekeeper calcula o score global a partir dos 7 "
             "pilares do OCG. Cada pilar tem peso canônico e mostra "
             "blockers e warnings de forma auditável.",
             size=10, color=INK_MED, line_spacing=1.4)

    # Box de regras
    bx = rx
    by = 4.05
    add_rect(slide, bx, by, rw, 2.4, fill=SURFACE_1, line=EDGE)
    add_text(slide, bx + 0.15, by + 0.12, rw - 0.3, 0.3,
             "REGRAS DE APROVAÇÃO",
             size=8, color=VIOLET_SOFT, bold=True)
    rules = [
        ("P7 < 70", "BLOCKED (segurança)", RED),
        ("P2 < 70", "BLOCKED (compliance)", RED),
        (">= 90", "READY", EMERALD),
        (">= 75", "NEEDS_REVIEW", AMBER),
        ("< 75", "AT_RISK", AMBER),
    ]
    ry = by + 0.45
    for cond, status, color in rules:
        add_text(slide, bx + 0.2, ry, rw * 0.35, 0.3,
                 cond, size=10, color=INK_MED, font="JetBrains Mono", bold=True)
        add_text(slide, bx + 0.2 + rw * 0.35, ry, rw * 0.6, 0.3,
                 f"→ {status}", size=10, color=color, bold=True)
        ry += 0.35

    add_footer(slide, idx, total)


def slide_07_pipeline(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "COMO FUNCIONA · 03",
        "Pipeline governado ponta a ponta.",
    )

    add_text(slide, MARGIN, 1.65, CONTENT_W, 0.4,
             "Cada etapa lê o OCG antes, atualiza o OCG depois. "
             "Auditoria SHA-256 em toda transição.",
             size=12, color=INK_LOW, italic=True, line_spacing=1.3)

    stages = [
        ("1", "Questionário", "Kickoff estruturado", VIOLET),
        ("2", "Repos + Ingestão", "Código legado + docs", VIOLET),
        ("3", "Gatekeeper", "7 pilares + score", VIOLET),
        ("4", "Arguidor", "Perguntas dirigidas", VIOLET),
        ("5", "OCG", "Consolidação", EMERALD),
        ("6", "CodeGen", "Scaffold + módulo", EMERALD),
        ("7", "Testes", "Unit + Int + E2E", EMERALD),
        ("8", "QA / Revisão", "Aprovação humana", EMERALD),
        ("9", "Roadmap", "Fases priorizadas", SKY),
        ("10", "Docs Vivas + ERS", "Sempre atualizado", SKY),
    ]

    cols = 5
    rows = 2
    gap_x, gap_y = 0.2, 0.25
    cw = (CONTENT_W - gap_x * (cols - 1)) / cols
    ch = 1.85
    start_y = 2.3

    for i, (num, title, sub, color) in enumerate(stages):
        r = i // cols
        c = i % cols
        cx = MARGIN + c * (cw + gap_x)
        cy = start_y + r * (ch + gap_y)
        add_rect(slide, cx, cy, cw, ch, fill=SURFACE_1, line=EDGE)
        # Badge numérico
        add_rect(slide, cx + 0.15, cy + 0.15, 0.5, 0.45,
                 fill=color, corner=0.35)
        add_text(slide, cx + 0.15, cy + 0.15, 0.5, 0.45,
                 num, size=12, color=INK_HIGH, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, cx + 0.15, cy + 0.7, cw - 0.3, 0.45,
                 title, size=12, color=INK_HIGH, bold=True)
        add_text(slide, cx + 0.15, cy + 1.15, cw - 0.3, 0.6,
                 sub, size=9, color=INK_LOW, line_spacing=1.25)

    # Callout final
    cy = start_y + 2 * ch + gap_y + 0.1
    add_text(slide, MARGIN, cy, CONTENT_W, 0.5,
             "Sem OCG, sem avanço. Sem aprovação humana, sem merge.",
             size=14, color=VIOLET_SOFT, bold=True, italic=True,
             align=PP_ALIGN.CENTER)

    add_footer(slide, idx, total)


def slide_08_codegen(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "COMO FUNCIONA · 04",
        "CodeGen — código que nasce dentro do contrato.",
    )

    # Screenshot à esquerda
    add_screenshot_frame(
        slide, MARGIN, 1.7, 8.3, 3.8,
        SHOTS / "codegen.png",
        caption="/projects/:id/codegen",
    )

    # Bloco de linguagens embaixo
    ly = 5.65
    add_text(slide, MARGIN, ly, CONTENT_W, 0.32,
             "8 LINGUAGENS · 10 SCAFFOLDERS ATUAIS · EXTENSÍVEL",
             size=10, color=VIOLET_SOFT, bold=True)
    langs = [
        "Python", "Node.js / Express", "Node.js / NestJS",
        "Java / Spring", "Java / Quarkus", "Kotlin / Spring",
        "C# / ASP.NET", "Go", "PHP / Laravel", "C++ / CMake",
    ]
    lx = MARGIN
    lyy = ly + 0.4
    for name in langs:
        w = add_pill(slide, lx, lyy, name, fill=SURFACE_2,
                     text_color=INK_HIGH, border=EDGE_STRONG, size=9)
        lx += w + 0.1
        if lx > MARGIN + CONTENT_W - 1:
            lx = MARGIN
            lyy += 0.42

    # Lado direito — explicação + Cobol/no-code
    rx = MARGIN + 8.3 + 0.25
    rw = CONTENT_W - 8.3 - 0.25

    add_text(slide, rx, 1.8, rw, 0.4,
             "Código, testes e docstring juntos.",
             size=13, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.25, rw, 1.2,
             "Cada arquivo gerado é rastreável ao requisito na "
             "matriz §4 do ERS e à trilha de auditoria SHA-256.",
             size=10, color=INK_MED, line_spacing=1.35)

    # Card "e amanhã?"
    ey = 3.6
    add_rect(slide, rx, ey, rw, 2.0, fill=SURFACE_3, line=VIOLET)
    add_text(slide, rx + 0.2, ey + 0.12, rw - 0.4, 0.35,
             "E AMANHÃ?",
             size=9, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx + 0.2, ey + 0.45, rw - 0.4, 0.45,
             "Expansível por scaffolder.",
             size=13, color=INK_HIGH, bold=True)
    add_text(slide, rx + 0.2, ey + 0.95, rw - 0.4, 1.0,
             "Cobol, Rust, Swift, Scala, RPG — ou no-code "
             "(Bubble, Retool, n8n) — entram como novos "
             "geradores. Nada na arquitetura impede.",
             size=10, color=INK_MED, line_spacing=1.35)

    add_footer(slide, idx, total)


def slide_09_roadmap(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "COMO FUNCIONA · 05",
        "Roadmap — o backlog em fases executáveis.",
    )

    add_screenshot_frame(
        slide, MARGIN, 1.75, 8.5, 4.9,
        SHOTS / "roadmap.png",
        caption="/projects/:id/roadmap",
    )

    rx = MARGIN + 8.5 + 0.25
    rw = CONTENT_W - 8.5 - 0.25

    add_text(slide, rx, 1.9, rw, 0.45,
             "Do questionário ao deploy.",
             size=14, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.4, rw, 1.2,
             "O Roadmap consolida módulos em fases "
             "(Fundação, Funcionalidades, Complementos) "
             "com filtros por categoria e status.",
             size=10, color=INK_MED, line_spacing=1.35)

    # 3 pontos
    items = [
        ("Fases", "ordem de deploy sugerida"),
        ("Filtros", "por tipo e prioridade"),
        ("Próxima ação", "recomendada pelo sistema"),
        ("Compartimentalizado", "por projeto (§2.2 canônico)"),
    ]
    iy = 3.75
    for t, sub in items:
        add_rect(slide, rx, iy, rw, 0.65, fill=SURFACE_1, line=EDGE)
        add_text(slide, rx + 0.2, iy + 0.08, rw - 0.4, 0.3,
                 t, size=11, color=INK_HIGH, bold=True)
        add_text(slide, rx + 0.2, iy + 0.35, rw - 0.4, 0.3,
                 sub, size=9, color=INK_LOW)
        iy += 0.7

    add_footer(slide, idx, total)


def slide_10_docs_vivas(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "COMO FUNCIONA · 06",
        "Documentação Viva — gerada a cada evento.",
    )

    add_screenshot_frame(
        slide, MARGIN, 1.75, 8.3, 4.9,
        SHOTS / "docs_vivas.png",
        caption="/projects/:id/livedocs",
    )

    rx = MARGIN + 8.3 + 0.25
    rw = CONTENT_W - 8.3 - 0.25

    add_text(slide, rx, 1.9, rw, 0.45,
             "Arquitetura · Índice · Módulos.",
             size=13, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.4, rw, 1.3,
             "Cada doc registra sua procedência: versão do OCG, "
             "ingestões usadas, provedor de IA e modelo.",
             size=10, color=INK_MED, line_spacing=1.35)

    # Destaque híbrido
    hy = 3.85
    add_rect(slide, rx, hy, rw, 2.6, fill=SURFACE_3, line=VIOLET)
    add_text(slide, rx + 0.2, hy + 0.15, rw - 0.4, 0.3,
             "ROTEAMENTO HÍBRIDO",
             size=8, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx + 0.2, hy + 0.5, rw - 0.4, 0.6,
             "Alto valor → premium\nLote repetitivo → local",
             size=11, color=INK_HIGH, bold=True, line_spacing=1.3)
    add_text(slide, rx + 0.2, hy + 1.35, rw - 0.4, 1.2,
             "Doc de Arquitetura usa Claude.\n"
             "Docs de 19 módulos rodam em Ollama local.\n"
             "Mesma entrega — fração do custo.",
             size=10, color=INK_MED, line_spacing=1.4)

    add_footer(slide, idx, total)


def slide_11_ers(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "COMO FUNCIONA · 07",
        "ERS vivo — IEEE 830 que não envelhece.",
    )

    # Mockup central em painel grande
    mx, my, mw, mh = MARGIN, 1.85, 7.8, 5.0
    add_rect(slide, mx, my, mw, mh, fill=SURFACE_1, line=EDGE)
    add_rect(slide, mx, my, mw, 0.4, fill=SURFACE_3, corner=0.12)
    add_text(slide, mx + 0.2, my + 0.08, mw - 0.4, 0.28,
             "docs/ERS.md  ·  commit e9a214f  ·  regenerado há 3 min",
             size=9, color=INK_LOW, font="JetBrains Mono")

    add_text(slide, mx + 0.3, my + 0.55, mw - 0.6, 0.45,
             "ERS — Especificação de Requisitos de Software",
             size=15, color=INK_HIGH, bold=True)
    add_pill(slide, mx + 0.3, my + 1.0, "IEEE 830-1998",
             fill=VIOLET_GLOW, text_color=VIOLET_SOFT,
             size=9, border=VIOLET, char_w=0.08)
    add_pill(slide, mx + 1.7, my + 1.0, "OCG v7",
             fill=SURFACE_3, text_color=INK_MED,
             size=9, border=EDGE_STRONG)
    add_pill(slide, mx + 2.5, my + 1.0, "✓ sincronizado",
             fill=SURFACE_3, text_color=EMERALD,
             size=9, border=EMERALD, char_w=0.08)

    sections = [
        ("1.3  Definições, Siglas e Abreviaturas", "12 termos aprovados"),
        ("3.1  Requisitos Funcionais (RF)", "27 requisitos"),
        ("3.2  Requisitos Não-Funcionais (RNF)", "14 requisitos"),
        ("3.3  Regras de Negócio (BR)", "9 regras"),
        ("4.0  Matriz de Rastreabilidade", "18 / 27 rastreados"),
    ]
    sy = my + 1.5
    for title, count in sections:
        add_rect(slide, mx + 0.3, sy, mw - 0.6, 0.45, fill=SURFACE_2, line=EDGE)
        add_text(slide, mx + 0.5, sy + 0.08, (mw - 1.0) * 0.62, 0.3,
                 title, size=10, color=INK_HIGH, bold=True)
        add_text(slide, mx + 0.5 + (mw - 1.0) * 0.62, sy + 0.08,
                 (mw - 1.0) * 0.38, 0.3,
                 count, size=10, color=VIOLET_SOFT, bold=True,
                 align=PP_ALIGN.RIGHT, font="JetBrains Mono")
        sy += 0.54

    # Git log mini
    gy = sy + 0.15
    add_rect(slide, mx + 0.3, gy, mw - 0.6, 0.95, fill=NAVY_DEEP, line=EDGE_STRONG)
    add_text(slide, mx + 0.45, gy + 0.08, mw - 0.9, 0.28,
             "$ git log --oneline docs/ERS.md",
             size=9, color=EMERALD, bold=True, font="JetBrains Mono")
    git_lines = [
        ("e9a214f", "docs(ers): regen OCG v7 — DOCUMENT_INGESTED"),
        ("7b3f41c", "docs(ers): regen OCG v6 — OCG_CONSOLIDATED"),
    ]
    ly = gy + 0.37
    for sha, msg in git_lines:
        add_text(slide, mx + 0.45, ly, 0.9, 0.25,
                 sha, size=9, color=AMBER, bold=True, font="JetBrains Mono")
        add_text(slide, mx + 1.45, ly, mw - 1.8, 0.25,
                 msg, size=9, color=INK_MED, font="JetBrains Mono")
        ly += 0.26

    # Coluna direita
    rx = MARGIN + mw + 0.3
    rw = CONTENT_W - mw - 0.3
    add_text(slide, rx, 1.9, rw, 0.4,
             "Sempre atualizado.",
             size=14, color=VIOLET_SOFT, bold=True)
    bullets = [
        "Regenerado no Git do projeto.",
        "Histórico nativo (git log).",
        "Zero snapshot em banco.",
        "Glossário aprovado pelo GP.",
        "Matriz §4 requisito × teste × código.",
        "Pandoc → PDF em um comando.",
    ]
    by = 2.5
    for b in bullets:
        add_text(slide, rx, by, 0.3, 0.32,
                 "✓", size=12, color=EMERALD, bold=True)
        add_text(slide, rx + 0.35, by, rw - 0.35, 0.45,
                 b, size=10, color=INK_MED, line_spacing=1.3)
        by += 0.45

    add_footer(slide, idx, total)


def slide_12_provedores(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "COMO FUNCIONA · 08",
        "Provedor de IA à sua escolha — com fallback local.",
    )

    add_screenshot_frame(
        slide, MARGIN, 1.75, 8.3, 4.2,
        SHOTS / "provedores.png",
        caption="/projects/:id/settings · Provedor de IA",
    )

    # Legenda abaixo
    ly = 6.05
    add_text(slide, MARGIN, ly, 8.3, 0.5,
             "Cliente escolhe: Anthropic, OpenAI, Google, Ollama local — "
             "multi-provedor por projeto + fallback hierárquico.",
             size=11, color=INK_LOW, italic=True, line_spacing=1.3)

    # Lado direito — proposta
    rx = MARGIN + 8.3 + 0.25
    rw = CONTENT_W - 8.3 - 0.25

    add_text(slide, rx, 1.85, rw, 0.45,
             "Nem refém, nem cego.",
             size=14, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.35, rw, 1.2,
             "Use a IA que melhor atende ao objetivo "
             "(qualidade, custo, latência, privacidade).",
             size=10, color=INK_MED, line_spacing=1.35)

    benefits = [
        ("Premium", "para OCG e decisão crítica", VIOLET),
        ("Local (Ollama)", "para lote repetitivo", EMERALD),
        ("Fallback", "se provedor cair, outro assume", SKY),
        ("Vault de chaves", "nunca em plaintext", AMBER),
    ]
    by = 3.75
    for t, sub, color in benefits:
        add_rect(slide, rx, by, rw, 0.75, fill=SURFACE_1, line=EDGE)
        add_rect(slide, rx, by, 0.1, 0.75, fill=color, corner=0)
        add_text(slide, rx + 0.25, by + 0.1, rw - 0.4, 0.3,
                 t, size=11, color=INK_HIGH, bold=True)
        add_text(slide, rx + 0.25, by + 0.4, rw - 0.4, 0.3,
                 sub, size=9, color=INK_LOW)
        by += 0.8

    add_footer(slide, idx, total)


def slide_13_custo(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "COMO FUNCIONA · 09",
        "Custo de IA visível — métricas por projeto.",
    )

    add_screenshot_frame(
        slide, MARGIN, 1.75, 8.5, 4.9,
        SHOTS / "metricas.png",
        caption="/projects/:id/metrics",
    )

    rx = MARGIN + 8.5 + 0.25
    rw = CONTENT_W - 8.5 - 0.25

    add_text(slide, rx, 1.9, rw, 0.45,
             "Sem surpresa na fatura.",
             size=14, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.4, rw, 1.15,
             "Cada chamada de LLM registra provedor, modelo, "
             "operação, tokens in/out e custo em USD.",
             size=10, color=INK_MED, line_spacing=1.35)

    # Box com highlights
    hy = 3.7
    add_rect(slide, rx, hy, rw, 2.8, fill=SURFACE_1, line=EDGE)
    add_text(slide, rx + 0.2, hy + 0.12, rw - 0.4, 0.3,
             "O GP VÊ",
             size=8, color=VIOLET_SOFT, bold=True)
    items = [
        ("Chamadas totais", "35"),
        ("Tokens in", "209.261"),
        ("Tokens out", "111.659"),
        ("Custo estimado", "$11.51"),
        ("Janela", "24h · 7d · 30d"),
    ]
    iy = hy + 0.45
    for label, value in items:
        add_text(slide, rx + 0.25, iy, rw * 0.55, 0.3,
                 label, size=10, color=INK_MED)
        add_text(slide, rx + 0.25 + rw * 0.55, iy, rw * 0.4, 0.3,
                 value, size=10, color=INK_HIGH, bold=True,
                 align=PP_ALIGN.RIGHT, font="JetBrains Mono")
        iy += 0.38

    add_text(slide, rx + 0.2, iy + 0.05, rw - 0.4, 0.3,
             "Breakdown por operação: analyzer, pillar_p1..p7, "
             "consolidator — GP decide onde otimizar.",
             size=8, color=INK_LOW, italic=True, line_spacing=1.3)

    add_footer(slide, idx, total)


def slide_14_backup(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "COMO FUNCIONA · 10",
        "Backup automático e rollback por projeto.",
    )

    add_screenshot_frame(
        slide, MARGIN, 1.75, 8.3, 4.4,
        SHOTS / "backup.png",
        caption="/projects/:id/backups",
    )

    rx = MARGIN + 8.3 + 0.25
    rw = CONTENT_W - 8.3 - 0.25

    add_text(slide, rx, 1.85, rw, 0.45,
             "Cron diário + sob demanda.",
             size=13, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx, 2.35, rw, 1.3,
             "Mantém os últimos 10 snapshots compartimentalizados "
             "por projeto. Rollback em um clique.",
             size=10, color=INK_MED, line_spacing=1.35)

    points = [
        ("Horário", "00 12 * * * (parametrizável)"),
        ("Retenção", "últimos 10 snapshots"),
        ("Integridade", "hash SHA por backup"),
        ("Manual", "GP ou Admin via botão"),
        ("Catch-up", "auto se servidor estava down"),
    ]
    py = 3.8
    for t, sub in points:
        add_rect(slide, rx, py, rw, 0.5, fill=SURFACE_1, line=EDGE)
        add_text(slide, rx + 0.2, py + 0.08, rw * 0.35, 0.3,
                 t, size=10, color=INK_HIGH, bold=True)
        add_text(slide, rx + 0.2 + rw * 0.35, py + 0.08, rw * 0.6, 0.3,
                 sub, size=9, color=INK_LOW)
        py += 0.55

    add_footer(slide, idx, total)


def slide_15_integracoes(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "COMO FUNCIONA · 11",
        "Integrações — o GCA conversa com suas ferramentas.",
    )

    # Hoje — lado esquerdo
    lx, ly, lw, lh = MARGIN, 1.85, 5.9, 4.8
    add_rect(slide, lx, ly, lw, lh, fill=SURFACE_1, line=EMERALD)
    add_rect(slide, lx, ly, lw, 0.1, fill=EMERALD, corner=0)
    add_text(slide, lx + 0.3, ly + 0.25, lw - 0.6, 0.4,
             "HOJE · disponível",
             size=10, color=EMERALD, bold=True)
    add_text(slide, lx + 0.3, ly + 0.7, lw - 0.6, 0.6,
             "Hospedagem de código",
             size=18, color=INK_HIGH, bold=True)

    repos = [
        ("GitHub · GitLab · Bitbucket", "commits, PRs/MRs, webhooks"),
        ("Jira · Trello", "backlog ↔ issue tracker (bidirecional)"),
        ("Sonar · Snyk · gitleaks", "findings alimentam P7 do OCG"),
        ("Slack", "eventos canônicos no canal do time"),
    ]
    ry = ly + 1.55
    for name, sub in repos:
        add_rect(slide, lx + 0.3, ry, lw - 0.6, 0.65, fill=SURFACE_2, line=EDGE)
        # badge com check
        ch = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(lx + 0.45), Inches(ry + 0.18),
            Inches(0.3), Inches(0.3),
        )
        _set_fill(ch, EMERALD)
        _set_line(ch, None)
        ch.shadow.inherit = False
        add_text(slide, lx + 0.45, ry + 0.18, 0.3, 0.3,
                 "✓", size=11, color=NAVY_DEEP, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, lx + 0.85, ry + 0.1, lw - 1.1, 0.3,
                 name, size=11, color=INK_HIGH, bold=True)
        add_text(slide, lx + 0.85, ry + 0.38, lw - 1.1, 0.3,
                 sub, size=9, color=INK_LOW)
        ry += 0.75

    # Amanhã — lado direito
    rx = MARGIN + lw + 0.25
    rw = CONTENT_W - lw - 0.25
    add_rect(slide, rx, ly, rw, lh, fill=SURFACE_1, line=VIOLET)
    add_rect(slide, rx, ly, rw, 0.1, fill=VIOLET, corner=0)
    add_text(slide, rx + 0.3, ly + 0.25, rw - 0.6, 0.4,
             "AMANHÃ · em evolução",
             size=10, color=VIOLET_SOFT, bold=True)
    add_text(slide, rx + 0.3, ly + 0.7, rw - 0.6, 0.6,
             "Fluxo, design e gestão",
             size=18, color=INK_HIGH, bold=True)

    future = [
        ("MS Teams · Mattermost", "notifier bi-direcional (ChatOps)", SKY),
        ("Linear · Asana · Monday", "trackers sob demanda (~1.5-2d cada)", SKY),
        ("Miro · Lucidchart", "diagramas colaborativos", SKY),
        ("Figma", "entrada UX/UI para CodeGen", ROSE),
    ]
    ry = ly + 1.55
    for name, sub, color in future:
        add_rect(slide, rx + 0.3, ry, rw - 0.6, 0.65, fill=SURFACE_2, line=EDGE)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(rx + 0.45), Inches(ry + 0.18),
            Inches(0.3), Inches(0.3),
        )
        _set_fill(dot, color)
        _set_line(dot, None)
        dot.shadow.inherit = False
        add_text(slide, rx + 0.45, ry + 0.18, 0.3, 0.3,
                 "~", size=12, color=NAVY_DEEP, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, rx + 0.85, ry + 0.1, rw - 1.1, 0.3,
                 name, size=12, color=INK_HIGH, bold=True)
        add_text(slide, rx + 0.85, ry + 0.36, rw - 1.1, 0.3,
                 sub, size=9, color=INK_LOW)
        ry += 0.75

    add_footer(slide, idx, total)


def slide_16_diferencial_cta(prs, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    paint_background(slide)
    add_header_accent(slide)
    add_title(
        slide, "POR QUE NÃO É MAIS UM COPILOT",
        "Copilot gera código. GCA gera produto governado.",
    )

    tools = ["Copilot", "Cursor", "v0/Lovable", "GCA"]
    rows_data = [
        ("Autocomplete de IA", True, True, True, True),
        ("Geração de módulo completo", False, "parcial", True, True),
        ("Contexto vivo do projeto (OCG)", False, False, False, True),
        ("ERS IEEE 830 auto-regenerado", False, False, False, True),
        ("Matriz requisito × teste × código", False, False, False, True),
        ("RBAC + auditoria SHA-256", False, False, False, True),
        ("LGPD / GDPR + quarentena PII", False, False, False, True),
        ("Multi-linguagem governada", False, "parcial", False, "8+"),
        ("Instalável on-premises", False, False, False, True),
        ("Controle de custo de IA", False, False, False, True),
    ]

    # Tabela — coordenadas seguras
    tx = MARGIN
    ty = 1.75
    tw = CONTENT_W
    row_h = 0.33
    cell_feat_w = tw * 0.42
    cell_tool_w = (tw - cell_feat_w) / len(tools)

    # Header
    add_rect(slide, tx, ty, tw, row_h, fill=SURFACE_2, line=EDGE)
    add_text(slide, tx + 0.15, ty, cell_feat_w - 0.15, row_h,
             "Recurso", size=10, color=INK_LOW, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    for i, t in enumerate(tools):
        cx = tx + cell_feat_w + i * cell_tool_w
        is_gca = t == "GCA"
        if is_gca:
            add_rect(slide, cx, ty, cell_tool_w, row_h,
                     fill=VIOLET_GLOW, line=VIOLET)
        add_text(slide, cx, ty, cell_tool_w, row_h,
                 t, size=10, color=VIOLET_SOFT if is_gca else INK_HIGH,
                 bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    ty += row_h

    # Rows
    for i, row in enumerate(rows_data):
        feat, *vals = row
        bg = SURFACE_1 if i % 2 == 0 else NAVY_DEEP
        add_rect(slide, tx, ty, tw, row_h, fill=bg, line=EDGE)
        add_text(slide, tx + 0.15, ty, cell_feat_w - 0.15, row_h,
                 feat, size=9, color=INK_MED, anchor=MSO_ANCHOR.MIDDLE)
        for j, v in enumerate(vals):
            is_gca = j == len(vals) - 1
            cx = tx + cell_feat_w + j * cell_tool_w
            if is_gca:
                add_rect(slide, cx, ty, cell_tool_w, row_h,
                         fill=VIOLET_GLOW, line=None)
            if v is True:
                text, color, size = "✓", EMERALD, 12
            elif v is False:
                text, color, size = "✗", RED, 12
            else:
                text, color, size = str(v), AMBER, 9
            add_text(slide, cx, ty, cell_tool_w, row_h,
                     text, size=size, color=color, bold=True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ty += row_h

    # Callout CTA
    cy = ty + 0.2
    cta_h = 1.15
    add_rect(slide, MARGIN, cy, CONTENT_W, cta_h, fill=SURFACE_1, line=VIOLET)
    add_logo_mark(slide, MARGIN + 0.3, cy + 0.25, w=1.7, h=0.65)

    add_text(slide, MARGIN + 2.2, cy + 0.18, CONTENT_W - 2.55, 0.45,
             "GCA — o primeiro produto que trata governança de IA como infraestrutura.",
             size=12, color=INK_HIGH, bold=True)
    add_text(slide, MARGIN + 2.2, cy + 0.6, CONTENT_W - 2.55, 0.5,
             "Instalável · on-premises · uma instância por cliente  →  fale comigo no privado.",
             size=10, color=VIOLET_SOFT, italic=True, line_spacing=1.25)

    add_footer(slide, idx, total)


# ─── Main ──────────────────────────────────────────────────────────────


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    builders = [
        slide_01_capa,
        slide_02_glossario,
        slide_03_dor,
        slide_04_nao_substitui,
        slide_05_ocg,
        slide_06_gatekeeper,
        slide_07_pipeline,
        slide_08_codegen,
        slide_09_roadmap,
        slide_10_docs_vivas,
        slide_11_ers,
        slide_12_provedores,
        slide_13_custo,
        slide_14_backup,
        slide_15_integracoes,
        slide_16_diferencial_cta,
    ]
    total = len(builders)
    for i, build in enumerate(builders, start=1):
        if build.__name__ == "slide_01_capa":
            build(prs, total)
        else:
            build(prs, i, total)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"✓ {OUT_PATH}  ({OUT_PATH.stat().st_size / 1024:.1f} KiB)")


if __name__ == "__main__":
    main()
