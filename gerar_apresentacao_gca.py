#!/usr/bin/env python3
"""
Gera apresentação comercial do GCA em PPTX
com marca visual: fundo escuro, violeta, fontes claras
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===================== CORES GCA =====================
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # #0F172A slate-900
BG_CARD = RGBColor(0x1E, 0x29, 0x3B)       # #1E293B slate-800
VIOLET = RGBColor(0x7C, 0x3A, 0xED)         # #7C3AED violet-600
VIOLET_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)   # #A78BFA violet-400
EMERALD = RGBColor(0x34, 0xD3, 0x99)        # #34D399 emerald-400
AMBER = RGBColor(0xFB, 0xBF, 0x24)          # #FBBF24 amber-400
RED = RGBColor(0xF8, 0x71, 0x71)            # #F87171 red-400
BLUE = RGBColor(0x60, 0xA5, 0xFA)           # #60A5FA blue-400
WHITE = RGBColor(0xF1, 0xF5, 0xF9)          # #F1F5F9 slate-100
GRAY = RGBColor(0x94, 0xA3, 0xB8)           # #94A3B8 slate-400
DARK_GRAY = RGBColor(0x64, 0x74, 0x8B)      # #64748B slate-500

LOGO_PATH = "/home/luiz/GCA/frontend/public/images/gca-logo-200.png"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_shape(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=GRAY, bullet_color=VIOLET_LIGHT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_before = Pt(6)
        # Bullet
        p.level = 0


def add_logo(slide, left=0.5, top=0.3, width=1.2):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(left), Inches(top), Inches(width))


# =====================================================================
# SLIDE 1: CAPA
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
add_logo(slide, left=5.5, top=1.0, width=2.3)
add_text(slide, 2.0, 3.5, 9.0, 1.0, "GCA", font_size=60, color=VIOLET_LIGHT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 2.0, 4.3, 9.0, 0.8, "Gestão de Codificação Assistida", font_size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 2.0, 5.2, 9.0, 0.6, "Governança inteligente para desenvolvimento de software", font_size=16, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 2.0, 6.5, 9.0, 0.4, "GCA Software  •  2026", font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 2: O PROBLEMA
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_text(slide, 0.5, 0.3, 8.0, 0.5, "O Problema", font_size=14, color=VIOLET_LIGHT, bold=True)
add_text(slide, 0.8, 1.2, 11.5, 0.8, "Projetos de software falham por falta de governança", font_size=32, color=WHITE, bold=True)

problems = [
    "70% dos projetos de TI excedem prazo e orçamento por requisitos mal definidos",
    "Documentação desatualizada gera retrabalho e código que não atende o negócio",
    "Decisões de arquitetura sem rastreabilidade comprometem qualidade e segurança",
    "Equipes gastam 40% do tempo em tarefas que poderiam ser automatizadas por IA",
    "Compliance (LGPD, GDPR) é verificado tarde demais — quando já é caro corrigir",
    "Sem visão unificada, GP não sabe o estado real do projeto até ser tarde",
]

for i, prob in enumerate(problems):
    y = 2.3 + i * 0.7
    add_shape(slide, 0.8, y, 0.08, 0.08, RED)
    add_text(slide, 1.2, y - 0.05, 11.0, 0.5, prob, font_size=15, color=GRAY)


# =====================================================================
# SLIDE 3: A SOLUÇÃO
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_text(slide, 0.5, 0.3, 8.0, 0.5, "A Solução", font_size=14, color=VIOLET_LIGHT, bold=True)
add_text(slide, 0.8, 1.2, 11.5, 0.8, "GCA — Inteligência Artificial aplicada à governança de código", font_size=28, color=WHITE, bold=True)

add_text(slide, 0.8, 2.3, 11.0, 1.0,
    "O GCA é uma plataforma que governa todo o ciclo de vida de um projeto de software: "
    "da concepção (questionário técnico) até a entrega (código gerado e testado), "
    "usando IA para analisar, validar e gerar artefatos com rastreabilidade total.",
    font_size=16, color=GRAY)

features = [
    ("🧠", "OCG — Inteligência Viva", "Objeto de Contexto Global que evolui com cada documento ingerido"),
    ("🔒", "7 Pilares de Qualidade", "Negócio, Compliance, Escopo, Performance, Arquitetura, Dados, Segurança"),
    ("🤖", "8 Agentes de IA", "Pipeline que avalia, classifica e consolida requisitos automaticamente"),
    ("📊", "Billing Transparente", "Custo de IA por projeto, por operação, em tempo real"),
]

for i, (icon, title, desc) in enumerate(features):
    x = 0.8 + (i % 2) * 6.0
    y = 3.8 + (i // 2) * 1.5
    add_shape(slide, x, y, 5.5, 1.2, BG_CARD)
    add_text(slide, x + 0.3, y + 0.15, 1.0, 0.5, icon, font_size=24)
    add_text(slide, x + 1.0, y + 0.15, 4.0, 0.4, title, font_size=16, color=WHITE, bold=True)
    add_text(slide, x + 1.0, y + 0.6, 4.0, 0.5, desc, font_size=12, color=GRAY)


# =====================================================================
# SLIDE 4: PIPELINE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_text(slide, 0.5, 0.3, 8.0, 0.5, "Pipeline do Projeto", font_size=14, color=VIOLET_LIGHT, bold=True)
add_text(slide, 0.8, 1.0, 11.5, 0.8, "Do questionário ao código — governado pelo OCG", font_size=28, color=WHITE, bold=True)

pipeline = [
    ("Questionário", "49 perguntas\n8 blocos", VIOLET),
    ("Repositório", "Git do projeto\n(bloqueante)", RED),
    ("Ingestão", "Docs, código\nPII detection", BLUE),
    ("Gatekeeper", "7 pilares\nscores", AMBER),
    ("Arguidor", "Gaps e lacunas\naltera OCG", EMERALD),
    ("Geração\nde Código", "IA + humano\ntestes auto", VIOLET),
    ("Backlog", "33+ itens\ndo OCG", BLUE),
    ("Docs Viva", "Auto-gerada\ndo OCG", EMERALD),
]

for i, (name, desc, color) in enumerate(pipeline):
    x = 0.5 + i * 1.55
    add_shape(slide, x, 2.2, 1.4, 2.0, BG_CARD)
    add_text(slide, x + 0.1, 2.3, 1.2, 0.6, name, font_size=11, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.1, 2.9, 1.2, 0.8, desc, font_size=9, color=GRAY, align=PP_ALIGN.CENTER)
    if i < len(pipeline) - 1:
        add_text(slide, x + 1.35, 2.9, 0.3, 0.5, "→", font_size=18, color=DARK_GRAY)

add_text(slide, 0.8, 4.6, 11.5, 0.6, "O OCG (Objeto de Contexto Global) é a fonte de verdade — expande com boa documentação, contrai com dados ruins.", font_size=14, color=GRAY)

# Detalhe dos 7 pilares
add_text(slide, 0.8, 5.3, 11.5, 0.5, "7 Pilares de Avaliação", font_size=18, color=WHITE, bold=True)
pilares = [
    ("P1 Negócio", "10%"), ("P2 Compliance", "15%"), ("P3 Escopo", "20%"),
    ("P4 Performance", "20%"), ("P5 Arquitetura", "15%"), ("P6 Dados", "10%"), ("P7 Segurança", "10%"),
]
for i, (name, peso) in enumerate(pilares):
    x = 0.8 + i * 1.7
    color = RED if "Compliance" in name or "Segurança" in name else VIOLET_LIGHT
    add_shape(slide, x, 5.9, 1.5, 0.8, BG_CARD)
    add_text(slide, x + 0.1, 5.95, 1.3, 0.35, name, font_size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.1, 6.3, 1.3, 0.3, peso, font_size=12, color=WHITE, align=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 5: OCG — INTELIGÊNCIA VIVA
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_text(slide, 0.5, 0.3, 8.0, 0.5, "OCG — Objeto de Contexto Global", font_size=14, color=VIOLET_LIGHT, bold=True)
add_text(slide, 0.8, 1.0, 11.5, 0.8, "A inteligência viva do seu projeto", font_size=28, color=WHITE, bold=True)

add_text(slide, 0.8, 2.0, 11.0, 1.0,
    "O OCG não é um documento estático. É um objeto de estado evolutivo orientado a eventos. "
    "Ele expande com boa ingestão de dados e contrai com dados ruins ou conflitantes. "
    "Nenhum módulo opera sem antes ler o OCG.",
    font_size=15, color=GRAY)

ocg_sections = [
    ("Perfil do Projeto", "Nome, tipo, criticidade, arquitetura"),
    ("Stack Recomendada", "Frontend, backend, banco, cache, infra"),
    ("Scores por Pilar", "7 pilares com pontuação e findings"),
    ("Compliance", "LGPD, GDPR, auditoria, PCI-DSS"),
    ("Estratégia de Testes", "Unitários, integração, E2E, segurança"),
    ("Análise de Riscos", "Alto, médio, baixo com mitigações"),
    ("Entregáveis", "Lista do que o projeto deve entregar"),
    ("Status de Aprovação", "READY, NEEDS_REVIEW, AT_RISK, BLOCKED"),
]

for i, (title, desc) in enumerate(ocg_sections):
    x = 0.8 + (i % 2) * 6.0
    y = 3.3 + (i // 2) * 0.9
    add_shape(slide, x, y, 5.5, 0.7, BG_CARD)
    add_text(slide, x + 0.3, y + 0.1, 2.5, 0.3, title, font_size=13, color=VIOLET_LIGHT, bold=True)
    add_text(slide, x + 0.3, y + 0.38, 4.8, 0.3, desc, font_size=11, color=GRAY)


# =====================================================================
# SLIDE 6: FUNCIONALIDADES
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_text(slide, 0.5, 0.3, 8.0, 0.5, "Funcionalidades", font_size=14, color=VIOLET_LIGHT, bold=True)
add_text(slide, 0.8, 1.0, 11.5, 0.6, "Tudo o que você precisa para governar projetos de software", font_size=26, color=WHITE, bold=True)

funcs = [
    ("Questionário Externo", "49 perguntas em 8 blocos com validação automática e score de aderência"),
    ("Repositório do Projeto", "Integração Git (GitHub/GitLab/Bitbucket) — bloqueante sem configuração"),
    ("Repositórios Externos", "Leitura read-only de repos de terceiros para enriquecer o OCG via IA"),
    ("Ingestão de Documentos", "Upload com detecção de PII, quarentena, hash SHA256, categorização"),
    ("Gatekeeper — 7 Pilares", "Avaliação automática com radar chart, findings e status bloqueante"),
    ("Arguidor Técnico", "Identifica gaps, show-stoppers e lacunas — solicita complementação"),
    ("Geração de Código", "IA gera código baseado no OCG, GP revisa e aprova, testes automáticos"),
    ("Backlog Vivo", "Derivado do OCG — regenera automaticamente quando o contexto muda"),
    ("Documentação Viva", "Gerada e atualizada automaticamente do OCG do projeto"),
    ("Billing por Projeto", "Custo de IA em USD por projeto, operação e provedor — transparência total"),
    ("Auditoria Hash Chain", "Trilha imutável com hash encadeado, correlation_id, 27 tipos de evento"),
    ("RBAC Compartimentalizado", "8 papéis, cada projeto isolado, chaves IA separadas Admin vs Projeto"),
]

for i, (title, desc) in enumerate(funcs):
    x = 0.5 + (i % 3) * 4.2
    y = 1.8 + (i // 3) * 1.3
    add_shape(slide, x, y, 3.9, 1.1, BG_CARD)
    add_text(slide, x + 0.2, y + 0.1, 3.5, 0.35, title, font_size=12, color=EMERALD, bold=True)
    add_text(slide, x + 0.2, y + 0.45, 3.5, 0.6, desc, font_size=10, color=GRAY)


# =====================================================================
# SLIDE 7: STACK TECNOLÓGICA
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_text(slide, 0.5, 0.3, 8.0, 0.5, "Stack Tecnológica", font_size=14, color=VIOLET_LIGHT, bold=True)
add_text(slide, 0.8, 1.0, 11.5, 0.6, "Construído com tecnologias modernas e escaláveis", font_size=26, color=WHITE, bold=True)

stack = [
    ("Backend", "Python 3.11 + FastAPI\n41 serviços, 24 routers\nSQLAlchemy async", VIOLET),
    ("Frontend", "React 18 + TypeScript\nVite + Tailwind CSS\n28 páginas", BLUE),
    ("Banco de Dados", "PostgreSQL 16\n50+ tabelas\nMigrations SQL", EMERALD),
    ("Cache & Mensageria", "Redis 7\nSessões, rate limiting\nLocks distribuídos", AMBER),
    ("IA Multi-Provider", "DeepSeek, Anthropic\nOpenAI, Grok, Gemini\n8 agentes paralelos", VIOLET_LIGHT),
    ("Automação", "n8n Workflows\n4 pipelines\nWebhooks", BLUE),
    ("Infraestrutura", "Docker Compose\n5 serviços\nCloudflare Tunnel", EMERALD),
    ("Segurança", "JWT RS256, RBAC\nVault criptografado\nHash chain auditoria", RED),
]

for i, (title, desc, color) in enumerate(stack):
    x = 0.5 + (i % 4) * 3.1
    y = 1.9 + (i // 4) * 2.5
    add_shape(slide, x, y, 2.9, 2.2, BG_CARD)
    add_text(slide, x + 0.2, y + 0.2, 2.5, 0.4, title, font_size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.2, y + 0.7, 2.5, 1.2, desc, font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 8: DIFERENCIAIS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_text(slide, 0.5, 0.3, 8.0, 0.5, "Diferenciais Competitivos", font_size=14, color=VIOLET_LIGHT, bold=True)
add_text(slide, 0.8, 1.0, 11.5, 0.6, "Por que o GCA é diferente", font_size=28, color=WHITE, bold=True)

diffs = [
    ("OCG Reativo", "O contexto do projeto evolui automaticamente com cada documento ingerido. Não é um relatório estático — é inteligência viva que expande ou contrai baseado na qualidade dos dados.", VIOLET_LIGHT),
    ("IA Agnóstica", "Suporte a 5 provedores de IA (DeepSeek, Anthropic, OpenAI, Grok, Gemini). O cliente escolhe o provider e modelo. Billing transparente por chamada.", BLUE),
    ("Compartimentalização Total", "Cada projeto é isolado: repositório próprio, chaves IA próprias, equipe própria, OCG próprio. Dados de um projeto nunca afetam outro.", EMERALD),
    ("Governança End-to-End", "Do questionário inicial ao código gerado e testado — cada etapa é rastreada, auditada e versionada com hash chain imutável.", AMBER),
    ("Código Aberto ao Humano", "A IA gera, mas o humano decide. Cada módulo de código passa por revisão do GP antes de ser commitado. IA + Humano = qualidade.", RED),
]

for i, (title, desc, color) in enumerate(diffs):
    y = 1.8 + i * 1.05
    add_shape(slide, 0.8, y, 11.5, 0.9, BG_CARD)
    add_text(slide, 1.1, y + 0.08, 3.0, 0.35, title, font_size=15, color=color, bold=True)
    add_text(slide, 1.1, y + 0.4, 10.8, 0.5, desc, font_size=11, color=GRAY)


# =====================================================================
# SLIDE 9: ENTREGÁVEIS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_text(slide, 0.5, 0.3, 8.0, 0.5, "Entregáveis", font_size=14, color=VIOLET_LIGHT, bold=True)
add_text(slide, 0.8, 1.0, 11.5, 0.6, "O que o cliente recebe ao usar o GCA", font_size=26, color=WHITE, bold=True)

deliverables = [
    "OCG completo do projeto com scores de 7 pilares e status de aprovação",
    "Análise de stack tecnológica recomendada com justificativa por camada",
    "Checklist de conformidade (LGPD, GDPR, PCI-DSS) com responsáveis",
    "Análise de riscos com mitigações e responsáveis atribuídos",
    "Backlog vivo com itens priorizados por categoria (módulos, testes, compliance)",
    "Código gerado por IA, revisado pelo GP, com testes unitários automáticos",
    "Documentação viva que se atualiza automaticamente com o OCG",
    "Trilha de auditoria completa com hash chain (rastreabilidade total)",
    "Billing detalhado de uso de IA por projeto (custo, tokens, operações)",
    "Dashboard executivo com radar de pilares, health do contexto e status",
]

for i, item in enumerate(deliverables):
    x = 0.8 + (i % 2) * 6.0
    y = 1.8 + (i // 2) * 0.95
    add_shape(slide, x, y, 0.08, 0.08, EMERALD)
    add_text(slide, x + 0.3, y - 0.08, 5.5, 0.5, item, font_size=13, color=GRAY)


# =====================================================================
# SLIDE 10: CONTATO
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide, left=5.5, top=1.5, width=2.3)
add_text(slide, 2.0, 3.5, 9.0, 0.8, "GCA", font_size=48, color=VIOLET_LIGHT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 2.0, 4.2, 9.0, 0.6, "Gestão de Codificação Assistida", font_size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 2.0, 5.0, 9.0, 0.5, "Governança inteligente. IA assistida. Código rastreável.", font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 2.0, 5.8, 9.0, 0.4, "gca.code-auditor.com.br", font_size=14, color=VIOLET_LIGHT, align=PP_ALIGN.CENTER)
add_text(slide, 2.0, 6.2, 9.0, 0.4, "GCA Software  •  Luiz Carlos Pielak", font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 2.0, 6.6, 9.0, 0.4, "pielak.ctba@gmail.com", font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# =====================================================================
# SALVAR
# =====================================================================
output = "/home/luiz/GCA/GCA_Apresentacao_Comercial.pptx"
prs.save(output)
print(f"✅ Apresentação gerada: {output}")
print(f"   10 slides, formato 16:9, marca visual GCA")
